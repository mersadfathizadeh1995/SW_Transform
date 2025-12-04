from __future__ import annotations

import os
import tkinter as tk
from tkinter import ttk, filedialog, messagebox

from sw_transform.core.service import run_single as svc_run_single, run_compare as svc_run_compare
from sw_transform.processing.registry import METHODS, dyn as dyn_func, compute_reverse_flag
from sw_transform.io.file_assignment import assign_files as assign_files_from_names

# Modular components
from sw_transform.gui.utils.defaults import DEFAULTS
from sw_transform.gui.utils.icons import load_icon, load_app_icon


class SimpleMASWGUI:
    def __init__(self, root: tk.Tk):
        self.root = root
        self.root.title("MASW – SW_Transform GUI")
        try:
            self.root.geometry("1000x680"); self.root.minsize(900, 560)
        except Exception:
            pass
        # App icon using utility function
        self._app_icon = load_app_icon(root)

        # icon cache
        self._icons: dict[str, tk.PhotoImage] = {}

        # state
        self.file_list: list[str] = []
        self.file_types: dict[str, str] = {}  # base -> 'seg2' or 'mat'
        self.offsets: dict[str, str] = {}
        self.reverse_flags: dict[str, bool] = {}
        self.output_folder: str = ""
        self.method_key = tk.StringVar(value="fk")
        self.vibrosis_mode = tk.BooleanVar(value=False)  # False=hammer, True=vibrosis
        self.cylindrical_var = tk.BooleanVar(value=False)  # Cylindrical steering for FDBF
        
        # Vibrosis .mat file array configuration
        self.dx_var = tk.StringVar(value="2.0")  # Sensor spacing (m) for .mat files
        
        # Auto-link: when vibrosis is checked, auto-check cylindrical
        self.vibrosis_mode.trace_add('write', self._on_vibrosis_changed)

        # inputs
        self.vmin_var = tk.StringVar(value="0"); self.vmax_var = tk.StringVar(value="5000")
        self.fmin_var = tk.StringVar(value="0"); self.fmax_var = tk.StringVar(value="100")
        self.time_start_var = tk.StringVar(value="0.0"); self.time_end_var = tk.StringVar(value="1.0")
        self.downsample_var = tk.BooleanVar(value=True); self.down_factor_var = tk.StringVar(value="16")
        self.numf_var = tk.StringVar(value="4000")
        self.grid_fk_var = tk.StringVar(value="4000"); self.tol_fk_var = tk.StringVar(value="0")
        self.grid_ps_var = tk.StringVar(value="1200"); self.vspace_ps_var = tk.StringVar(value="log"); self.tol_ps_var = tk.StringVar(value="0")
        self.dpi_var = tk.StringVar(value="200")
        self.figure_topic_var = tk.StringVar(value="3-D Dispersion (Freq vs. Velocity)")
        self.ppt_var = tk.BooleanVar(value=False)
        self.export_spectra_var = tk.BooleanVar(value=True)  # Default ON for spectrum export
        self.parallel_var = tk.BooleanVar(value=True)  # Enable parallel processing by default
        self.worker_count_var = tk.StringVar(value="auto")  # "auto" or specific number
        
        # === Advanced Settings Variables (use imported DEFAULTS) ===
        # New variables for advanced settings
        self.power_threshold_var = tk.StringVar(value=DEFAULTS['power_threshold'])
        self.auto_vel_limits_var = tk.BooleanVar(value=DEFAULTS['auto_vel_limits'])
        self.auto_freq_limits_var = tk.BooleanVar(value=DEFAULTS['auto_freq_limits'])
        self.plot_min_vel_var = tk.StringVar(value=DEFAULTS['plot_min_vel'])
        self.plot_max_vel_var = tk.StringVar(value=DEFAULTS['plot_max_vel'])
        self.plot_min_freq_var = tk.StringVar(value=DEFAULTS['plot_min_freq'])
        self.plot_max_freq_var = tk.StringVar(value=DEFAULTS['plot_max_freq'])
        self.freq_tick_spacing_var = tk.StringVar(value=DEFAULTS['freq_tick_spacing'])
        self.vel_tick_spacing_var = tk.StringVar(value=DEFAULTS['vel_tick_spacing'])
        self.cmap_var = tk.StringVar(value=DEFAULTS['cmap'])

        self._build_menu()
        self._build_ui()

    # UI
    def _build_menu(self):
        m = tk.Menu(self.root)
        filem = tk.Menu(m, tearoff=0)
        filem.add_command(label="Open Data Files...", command=self.select_files)
        filem.add_command(label="Open Vibrosis .MAT...", command=self.select_mat_files)
        filem.add_command(label="Select Output Folder...", command=self.select_out)
        filem.add_separator(); filem.add_command(label="Exit", command=self.root.quit)
        m.add_cascade(label="File", menu=filem)

        datam = tk.Menu(m, tearoff=0)
        datam.add_command(label="Auto Assign Offsets (from filenames)", command=self._auto_assign_from_filenames)
        m.add_cascade(label="Data", menu=datam)

        self.root.config(menu=m)

    def _build_ui(self):
        left = tk.Frame(self.root, width=320); left.pack(side="left", fill="y")
        # Left toolbar with icon
        btn_open = tk.Button(left, text=" Open Data...", command=self.select_files, compound="left", padx=6, pady=4)
        ico = self._load_icon("ic_open.png", 32)
        if ico is not None:
            btn_open.config(image=ico)
        btn_open.pack(anchor="w", padx=8, pady=6)
        self.tree = ttk.Treeview(left, columns=("file","type","offset","rev"), show="headings", height=24)
        for col, w in zip(("file","type","offset","rev"), (180, 40, 60, 40)):
            self.tree.heading(col, text=col.capitalize()); self.tree.column(col, width=w)
        self.tree.pack(fill="y", padx=8, pady=4, expand=True)
        self.tree.bind("<Double-1>", self._edit_cell)

        center = tk.Frame(self.root); center.pack(side="left", fill="both", expand=True)
        nb = ttk.Notebook(center)
        self.tab_inputs = tk.Frame(nb); self.tab_run = tk.Frame(nb); self.tab_fig = tk.Frame(nb)
        nb.add(self.tab_inputs, text="Inputs"); nb.add(self.tab_run, text="Run"); nb.add(self.tab_fig, text="Figures")
        
        # Add MASW 2D tab
        try:
            from sw_transform.gui.masw2d_tab import create_masw2d_tab
            self.masw2d_tab = create_masw2d_tab(nb, log_callback=None, main_app=self)
        except ImportError:
            pass  # MASW 2D module not available
        
        nb.pack(fill="both", expand=True)

        # Inputs
        p = self.tab_inputs
        row = tk.Frame(p); row.pack(fill="x", padx=6, pady=4)
        tk.Label(row, text="Output folder:").pack(side="left")
        self.out_var = tk.StringVar(value="(not set)")
        tk.Label(row, textvariable=self.out_var, anchor="w").pack(side="left", fill="x", expand=True, padx=6)
        tk.Button(row, text="Select", command=self.select_out).pack(side="left")

        # Processing Limits (simplified)
        limits_box = tk.LabelFrame(p, text="Processing Limits"); limits_box.pack(fill="x", padx=6, pady=4)
        lim_row1 = tk.Frame(limits_box); lim_row1.pack(fill="x", pady=2)
        tk.Label(lim_row1, text="Velocity:").pack(side="left")
        tk.Entry(lim_row1, width=6, textvariable=self.vmin_var).pack(side="left", padx=2)
        tk.Label(lim_row1, text="-").pack(side="left")
        tk.Entry(lim_row1, width=6, textvariable=self.vmax_var).pack(side="left", padx=2)
        tk.Label(lim_row1, text="m/s").pack(side="left", padx=(0, 12))
        tk.Label(lim_row1, text="Frequency:").pack(side="left")
        tk.Entry(lim_row1, width=6, textvariable=self.fmin_var).pack(side="left", padx=2)
        tk.Label(lim_row1, text="-").pack(side="left")
        tk.Entry(lim_row1, width=6, textvariable=self.fmax_var).pack(side="left", padx=2)
        tk.Label(lim_row1, text="Hz").pack(side="left")
        
        lim_row2 = tk.Frame(limits_box); lim_row2.pack(fill="x", pady=2)
        tk.Label(lim_row2, text="Time Window:").pack(side="left")
        tk.Entry(lim_row2, width=6, textvariable=self.time_start_var).pack(side="left", padx=2)
        tk.Label(lim_row2, text="-").pack(side="left")
        tk.Entry(lim_row2, width=6, textvariable=self.time_end_var).pack(side="left", padx=2)
        tk.Label(lim_row2, text="sec").pack(side="left")

        # Topic / Title
        topic_box = tk.LabelFrame(p, text="Figure Title"); topic_box.pack(fill="x", padx=6, pady=4)
        tk.Entry(topic_box, textvariable=self.figure_topic_var, width=50).pack(fill="x", padx=4, pady=4)

        # Advanced Settings button
        adv_row = tk.Frame(p); adv_row.pack(fill="x", padx=6, pady=4)
        tk.Button(adv_row, text="\u2699 Advanced Settings...", command=self._open_advanced_settings).pack(side="left")

        # Array preview (embedded)
        arr_box = tk.LabelFrame(p, text="Array preview (embedded)")
        arr_box.pack(fill="both", expand=True, padx=6, pady=6)
        topbar = tk.Frame(arr_box); topbar.pack(fill="x", pady=(2,4))
        tk.Button(topbar, text="Preview Array / Waterfall", command=self.preview_array).pack(side="left")
        tk.Label(topbar, text="Display time (s):").pack(side="left", padx=(10,2))
        self.display_time_var = tk.StringVar(value="1")
        tk.Entry(topbar, width=6, textvariable=self.display_time_var).pack(side="left")
        self.prev_host = tk.Frame(arr_box, height=300, bg="#f7f7f7"); self.prev_host.pack(fill="both", expand=True)
        self.prev_canvas_widget = None

        # Run tab
        r = self.tab_run
        top = tk.Frame(r); top.pack(pady=6)
        tk.Label(top, text="Transform:").pack(side="left")
        cmb = ttk.Combobox(top, values=[f"{k} – {d['label']}" for k,d in METHODS.items()], state="readonly", width=48)
        cmb.current(0); cmb.pack(side="left", padx=6)
        cmb.bind("<<ComboboxSelected>>", lambda ev: self.method_key.set(ev.widget.get().split(" –")[0]))
        row = tk.Frame(r); row.pack(pady=4)
        btn_run_sel = tk.Button(row, text=" Run Selected", command=lambda: self.run_single_processing(selected_only=True), compound="left", padx=8, pady=4)
        btn_run_all = tk.Button(row, text=" Run All", command=lambda: self.run_single_processing(selected_only=False), compound="left", padx=8, pady=4)
        ico_run = self._load_icon("ic_run.png", 32)
        if ico_run is not None:
            btn_run_sel.config(image=ico_run); btn_run_all.config(image=ico_run)
        btn_run_sel.pack(side="left", padx=4); btn_run_all.pack(side="left", padx=4)
        row2 = tk.Frame(r); row2.pack(pady=4)
        btn_cmp_sel = tk.Button(row2, text=" Compare Selected", command=lambda: self.run_compare_processing(selected_only=True), compound="left", padx=8, pady=4)
        btn_cmp_all = tk.Button(row2, text=" Compare All", command=lambda: self.run_compare_processing(selected_only=False), compound="left", padx=8, pady=4)
        ico_cmp = self._load_icon("ic_compare.png", 32)
        if ico_cmp is not None:
            btn_cmp_sel.config(image=ico_cmp); btn_cmp_all.config(image=ico_cmp)
        btn_cmp_sel.pack(side="left", padx=4); btn_cmp_all.pack(side="left", padx=4)
        opt = tk.Frame(r); opt.pack(pady=2)
        tk.Checkbutton(opt, text="Create PowerPoint after run", variable=self.ppt_var).pack(side="left")
        tk.Checkbutton(opt, text="Parallel processing", variable=self.parallel_var).pack(side="left", padx=(16, 0))
        # Worker count control
        tk.Label(opt, text="Workers:").pack(side="left", padx=(8, 2))
        import multiprocessing
        max_cpu = multiprocessing.cpu_count()
        worker_combo = ttk.Combobox(opt, textvariable=self.worker_count_var, 
                                     values=["auto"] + [str(i) for i in range(1, max_cpu + 1)],
                                     width=5, state="readonly")
        worker_combo.pack(side="left")
        tk.Label(opt, text=f"(max {max_cpu})", fg="gray").pack(side="left", padx=(2, 0))
        # Log box
        try:
            from tkinter import scrolledtext
            self.logbox = scrolledtext.ScrolledText(r, width=92, height=10)
            self.logbox.pack(fill="both", expand=True, padx=6, pady=6)
        except Exception:
            self.logbox = None

        # Bottom-right progress bar
        pb_row = tk.Frame(r); pb_row.pack(fill="x", padx=6, pady=(0,6))
        self.pb = ttk.Progressbar(pb_row, orient="horizontal", mode="determinate", length=220)
        self.pb.pack(side="right")
        self.pb_label = tk.Label(pb_row, text="Idle", anchor="e")
        self.pb_label.pack(side="right", padx=(0,8))

        # Figures tab UI (split: center preview + right sidebar explorer)
        ftab = self.tab_fig
        outer = tk.Frame(ftab); outer.pack(fill="both", expand=True)
        # Center preview area
        center = tk.Frame(outer); center.pack(side="left", fill="both", expand=True, padx=(6,0), pady=6)
        toolbar = tk.Frame(center); toolbar.pack(fill="x")
        tk.Button(toolbar, text="−", width=3, command=lambda: self._fig_zoom_step(1/1.15)).pack(side="left")
        tk.Button(toolbar, text="+", width=3, command=lambda: self._fig_zoom_step(1.15)).pack(side="left", padx=(4,0))
        tk.Button(toolbar, text="100%", command=self._fig_zoom_reset).pack(side="left", padx=(8,0))
        tk.Label(toolbar, text="Fit:").pack(side="left", padx=(12,2))
        self.fig_fit_mode = tk.StringVar(value="Auto")
        ttk.Combobox(toolbar, values=("Auto","Width","Height","None"), width=8, state="readonly", textvariable=self.fig_fit_mode).pack(side="left")
        self.fig_zoom_label = tk.StringVar(value="100%")
        tk.Label(toolbar, textvariable=self.fig_zoom_label).pack(side="left", padx=(10,0))
        # Preview canvas
        prev_frame = tk.Frame(center); prev_frame.pack(fill="both", expand=True, pady=(4,0))
        self.fig_prev_canvas = tk.Canvas(prev_frame, bg="#f3f3f3")
        vs = tk.Scrollbar(prev_frame, orient="vertical", command=self.fig_prev_canvas.yview)
        hs = tk.Scrollbar(prev_frame, orient="horizontal", command=self.fig_prev_canvas.xview)
        self.fig_prev_canvas.configure(yscrollcommand=vs.set, xscrollcommand=hs.set)
        self.fig_prev_canvas.grid(row=0, column=0, sticky="nsew"); vs.grid(row=0, column=1, sticky="ns"); hs.grid(row=1, column=0, sticky="ew")
        prev_frame.rowconfigure(0, weight=1); prev_frame.columnconfigure(0, weight=1)
        self.fig_prev_canvas.bind("<ButtonPress-1>", lambda e: self.fig_prev_canvas.scan_mark(e.x, e.y))
        self.fig_prev_canvas.bind("<B1-Motion>", lambda e: self.fig_prev_canvas.scan_dragto(e.x, e.y, gain=1))
        self.fig_prev_canvas.bind("<MouseWheel>", self._on_fig_wheel_zoom)
        self.fig_prev_canvas.bind("<Configure>", lambda e: self._render_fig_preview())

        # Right sidebar explorer
        right = tk.Frame(outer, width=320); right.pack(side="left", fill="y", padx=6, pady=6)
        sb_tools = tk.Frame(right); sb_tools.pack(fill="x")
        tk.Button(sb_tools, text="Refresh", command=self.refresh_gallery, pady=2).pack(side="left")
        tk.Button(sb_tools, text="Open", command=self._open_selected_figure, pady=2).pack(side="left", padx=(6,0))
        tk.Button(sb_tools, text="Delete", command=self._delete_selected_figure, pady=2).pack(side="left", padx=(6,0))
        btn_ppt = tk.Button(sb_tools, text=" PPT", command=self._build_ppt_from_gallery, compound="left", padx=6, pady=2)
        ico_ppt = self._load_icon("ic_ppt.png", 28)
        if ico_ppt is not None:
            btn_ppt.config(image=ico_ppt)
        btn_ppt.pack(side="left", padx=(6,0))
        self.fig_list = ttk.Treeview(right, columns=("name","path"), show="headings", height=22)
        self.fig_list.heading("name", text="Name"); self.fig_list.column("name", width=180)
        self.fig_list.heading("path", text="Path"); self.fig_list.column("path", width=120)
        self.fig_list.pack(fill="both", expand=True, pady=(6,0))
        self.fig_list.bind("<<TreeviewSelect>>", self._on_figure_selected)

        # Preview state
        self._fig_image_pil = None
        self._fig_image_tk = None
        self._fig_scale = 1.0

    # actions
    def select_files(self):
        """Open file dialog for SEG-2 (.dat) and vibrosis (.mat) files."""
        files = filedialog.askopenfilenames(
            title="Select Data Files",
            filetypes=[
                ("All supported", "*.dat *.mat"),
                ("SEG-2 .dat", "*.dat"),
                ("Vibrosis .mat", "*.mat"),
                ("All files", "*.*")
            ]
        )
        if not files:
            return
        self._add_files_to_list(files)
    
    def select_mat_files(self):
        """Open file dialog specifically for vibrosis .mat files."""
        files = filedialog.askopenfilenames(
            title="Select Vibrosis .MAT Files",
            filetypes=[("Vibrosis .mat", "*.mat"), ("All files", "*.*")]
        )
        if not files:
            return
        self._add_files_to_list(files)
    
    def _add_files_to_list(self, files: tuple | list):
        """Add files to the file list, detecting type and auto-configuring.
        
        Supports mixed mode: both .dat (SEG-2) and .mat (vibrosis) files.
        """
        new_files = list(files)
        
        # Extend existing lists (mixed mode)
        self.file_list.extend(new_files)
        
        # Process each new file
        for f in new_files:
            base = os.path.splitext(os.path.basename(f))[0]
            ext = os.path.splitext(f)[1].lower()
            
            if ext == '.mat':
                # Vibrosis .mat file
                self.file_types[base] = 'mat'
                self.reverse_flags[base] = False
                
                # Try to detect array info and parse offset
                try:
                    from sw_transform.processing.vibrosis import get_vibrosis_file_info
                    info = get_vibrosis_file_info(f)
                    parsed_offset = info.get('parsed_offset')
                    if parsed_offset is not None:
                        sign = "+" if parsed_offset >= 0 else ""
                        self.offsets[base] = f"{sign}{int(round(parsed_offset))}"
                    else:
                        self.offsets[base] = "+0"
                except Exception:
                    self.offsets[base] = "+0"
            else:
                # SEG-2 .dat file
                self.file_types[base] = 'seg2'
                self.offsets[base] = "+0"
                self.reverse_flags[base] = False
        
        # Refresh tree view
        self._refresh_file_tree()
        
        # Auto-enable vibrosis mode if any .mat files loaded
        has_mat = any(t == 'mat' for t in self.file_types.values())
        if has_mat:
            self.vibrosis_mode.set(True)
    
    def _refresh_file_tree(self):
        """Refresh the file tree view with current file list."""
        self.tree.delete(*self.tree.get_children())
        for f in self.file_list:
            base = os.path.splitext(os.path.basename(f))[0]
            ftype = self.file_types.get(base, 'seg2')
            offset = self.offsets.get(base, '+0')
            rev = self.reverse_flags.get(base, False)
            type_label = "MAT" if ftype == 'mat' else "SEG2"
            rev_label = "☑" if rev else "☐"
            self.tree.insert("", "end", values=(base, type_label, offset, rev_label))

    def select_out(self):
        folder = filedialog.askdirectory(title="Choose output folder")
        if folder:
            self.output_folder = folder; self.out_var.set(folder)

    def _on_vibrosis_changed(self, *args):
        """Auto-check cylindrical when vibrosis is enabled."""
        if self.vibrosis_mode.get():
            self.cylindrical_var.set(True)
    
    def _reset_advanced_defaults(self):
        """Reset all advanced settings to default values."""
        self.grid_fk_var.set(DEFAULTS['grid_fk'])
        self.tol_fk_var.set(DEFAULTS['tol_fk'])
        self.grid_ps_var.set(DEFAULTS['grid_ps'])
        self.vspace_ps_var.set(DEFAULTS['vspace_ps'])
        self.tol_ps_var.set(DEFAULTS['tol_ps'])
        self.vibrosis_mode.set(DEFAULTS['vibrosis'])
        self.cylindrical_var.set(DEFAULTS['cylindrical'])
        self.downsample_var.set(DEFAULTS['downsample'])
        self.down_factor_var.set(DEFAULTS['down_factor'])
        self.numf_var.set(DEFAULTS['numf'])
        self.power_threshold_var.set(DEFAULTS['power_threshold'])
        self.auto_vel_limits_var.set(DEFAULTS['auto_vel_limits'])
        self.auto_freq_limits_var.set(DEFAULTS['auto_freq_limits'])
        self.plot_min_vel_var.set(DEFAULTS['plot_min_vel'])
        self.plot_max_vel_var.set(DEFAULTS['plot_max_vel'])
        self.plot_min_freq_var.set(DEFAULTS['plot_min_freq'])
        self.plot_max_freq_var.set(DEFAULTS['plot_max_freq'])
        self.freq_tick_spacing_var.set(DEFAULTS['freq_tick_spacing'])
        self.vel_tick_spacing_var.set(DEFAULTS['vel_tick_spacing'])
        self.cmap_var.set(DEFAULTS['cmap'])
        self.dpi_var.set(DEFAULTS['dpi'])
        self.export_spectra_var.set(DEFAULTS['export_spectra'])
        self.dx_var.set(DEFAULTS['dx'])
        self._toggle_vel_limits()
        self._toggle_freq_limits()
    
    def _toggle_vel_limits(self):
        """Enable/disable velocity limit entries based on auto-limit checkbox."""
        state = 'disabled' if self.auto_vel_limits_var.get() else 'normal'
        if hasattr(self, 'plot_min_vel_entry'):
            self.plot_min_vel_entry.config(state=state)
            self.plot_max_vel_entry.config(state=state)
    
    def _toggle_freq_limits(self):
        """Enable/disable frequency limit entries based on auto-limit checkbox."""
        state = 'disabled' if self.auto_freq_limits_var.get() else 'normal'
        if hasattr(self, 'plot_min_freq_entry'):
            self.plot_min_freq_entry.config(state=state)
            self.plot_max_freq_entry.config(state=state)

    def _open_advanced_settings(self):
        """Open the advanced settings popup window."""
        popup = tk.Toplevel(self.root)
        popup.title("Advanced Settings")
        popup.geometry("480x520")
        popup.resizable(True, True)
        popup.transient(self.root)
        popup.grab_set()
        
        # Main scrollable frame
        canvas = tk.Canvas(popup, highlightthickness=0)
        scrollbar = ttk.Scrollbar(popup, orient="vertical", command=canvas.yview)
        scrollable = tk.Frame(canvas)
        
        scrollable.bind("<Configure>", lambda e: canvas.configure(scrollregion=canvas.bbox("all")))
        canvas.create_window((0, 0), window=scrollable, anchor="nw")
        canvas.configure(yscrollcommand=scrollbar.set)
        
        scrollbar.pack(side="right", fill="y")
        canvas.pack(side="left", fill="both", expand=True, padx=5, pady=5)
        
        # ===== Transform Settings =====
        transform_frame = tk.LabelFrame(scrollable, text="Transform Settings", padx=8, pady=8)
        transform_frame.pack(fill="x", padx=8, pady=6)
        
        # FK/FDBF row
        fk_row = tk.Frame(transform_frame); fk_row.pack(fill="x", pady=3)
        tk.Label(fk_row, text="FK/FDBF Grid Size:", width=16, anchor="w").pack(side="left")
        ttk.Combobox(fk_row, textvariable=self.grid_fk_var,
                     values=["500", "1000", "2000", "4000", "8000"],
                     width=8).pack(side="left", padx=4)
        tk.Label(fk_row, text="Tolerance:").pack(side="left", padx=(12, 0))
        tk.Entry(fk_row, textvariable=self.tol_fk_var, width=8).pack(side="left", padx=4)
        
        # PS/SS row
        ps_row = tk.Frame(transform_frame); ps_row.pack(fill="x", pady=3)
        tk.Label(ps_row, text="PS/SS Grid Size:", width=16, anchor="w").pack(side="left")
        ttk.Combobox(ps_row, textvariable=self.grid_ps_var,
                     values=["500", "1000", "1200", "2000", "4000", "8000"],
                     width=8).pack(side="left", padx=4)
        tk.Label(ps_row, text="Tolerance:").pack(side="left", padx=(12, 0))
        tk.Entry(ps_row, textvariable=self.tol_ps_var, width=8).pack(side="left", padx=4)
        
        # PS/SS spacing row
        sp_row = tk.Frame(transform_frame); sp_row.pack(fill="x", pady=3)
        tk.Label(sp_row, text="PS/SS Spacing:", width=16, anchor="w").pack(side="left")
        ttk.Combobox(sp_row, textvariable=self.vspace_ps_var,
                     values=["log", "linear"], width=8, state="readonly").pack(side="left", padx=4)
        
        # Vibrosis mode
        vib_row = tk.Frame(transform_frame); vib_row.pack(fill="x", pady=3)
        tk.Checkbutton(vib_row, text="Vibrosis mode (FDBF weighting)", 
                       variable=self.vibrosis_mode).pack(side="left")
        
        # Cylindrical steering
        cyl_row = tk.Frame(transform_frame); cyl_row.pack(fill="x", pady=3)
        tk.Checkbutton(cyl_row, text="Cylindrical steering (FDBF near-field)", 
                       variable=self.cylindrical_var).pack(side="left")
        
        # ===== Vibrosis Array Settings =====
        vib_array_frame = tk.LabelFrame(scrollable, text="Vibrosis Array Config (.mat files)", padx=8, pady=8)
        vib_array_frame.pack(fill="x", padx=8, pady=6)
        
        # Sensor spacing (dx)
        dx_row = tk.Frame(vib_array_frame); dx_row.pack(fill="x", pady=3)
        tk.Label(dx_row, text="Sensor Spacing (dx):", width=16, anchor="w").pack(side="left")
        tk.Entry(dx_row, textvariable=self.dx_var, width=10).pack(side="left", padx=4)
        tk.Label(dx_row, text="meters", fg="gray").pack(side="left")
        
        # Info label
        tk.Label(vib_array_frame, text="Channel count auto-detected from .mat files", 
                 fg="gray", font=("TkDefaultFont", 8)).pack(anchor="w", pady=(0, 3))
        
        # ===== Preprocessing Settings =====
        preproc_frame = tk.LabelFrame(scrollable, text="Preprocessing", padx=8, pady=8)
        preproc_frame.pack(fill="x", padx=8, pady=6)
        
        # Downsample row
        ds_row = tk.Frame(preproc_frame); ds_row.pack(fill="x", pady=3)
        tk.Checkbutton(ds_row, text="Downsample", variable=self.downsample_var).pack(side="left")
        tk.Label(ds_row, text="Factor:").pack(side="left", padx=(16, 0))
        ttk.Combobox(ds_row, textvariable=self.down_factor_var,
                     values=["1", "2", "4", "8", "16", "32"],
                     width=6, state="readonly").pack(side="left", padx=4)
        
        # FFT size row
        fft_row = tk.Frame(preproc_frame); fft_row.pack(fill="x", pady=3)
        tk.Label(fft_row, text="FFT Size (numf):", width=16, anchor="w").pack(side="left")
        ttk.Combobox(fft_row, textvariable=self.numf_var,
                     values=["1000", "2000", "4000", "8000"],
                     width=8).pack(side="left", padx=4)
        tk.Label(fft_row, text="points", fg="gray").pack(side="left")
        
        # ===== Peak Picking Settings =====
        peak_frame = tk.LabelFrame(scrollable, text="Peak Picking", padx=8, pady=8)
        peak_frame.pack(fill="x", padx=8, pady=6)
        
        pt_row = tk.Frame(peak_frame); pt_row.pack(fill="x", pady=3)
        tk.Label(pt_row, text="Power Threshold:", width=16, anchor="w").pack(side="left")
        tk.Entry(pt_row, textvariable=self.power_threshold_var, width=10).pack(side="left", padx=4)
        tk.Label(pt_row, text="(0.0-1.0)", fg="gray").pack(side="left")
        
        # ===== Image Export Settings =====
        image_frame = tk.LabelFrame(scrollable, text="Image Export", padx=8, pady=8)
        image_frame.pack(fill="x", padx=8, pady=6)
        
        # Velocity range for plots with auto checkbox
        vel_row = tk.Frame(image_frame); vel_row.pack(fill="x", pady=3)
        tk.Checkbutton(vel_row, text="Auto", variable=self.auto_vel_limits_var,
                       command=self._toggle_vel_limits).pack(side="left")
        tk.Label(vel_row, text="Velocity:", width=8, anchor="w").pack(side="left")
        self.plot_min_vel_entry = tk.Entry(vel_row, textvariable=self.plot_min_vel_var, width=7)
        self.plot_min_vel_entry.pack(side="left", padx=2)
        tk.Label(vel_row, text="to").pack(side="left", padx=2)
        self.plot_max_vel_entry = tk.Entry(vel_row, textvariable=self.plot_max_vel_var, width=7)
        self.plot_max_vel_entry.pack(side="left", padx=2)
        tk.Label(vel_row, text="m/s", fg="gray").pack(side="left", padx=2)
        
        # Frequency range for plots with auto checkbox
        freq_row = tk.Frame(image_frame); freq_row.pack(fill="x", pady=3)
        tk.Checkbutton(freq_row, text="Auto", variable=self.auto_freq_limits_var,
                       command=self._toggle_freq_limits).pack(side="left")
        tk.Label(freq_row, text="Frequency:", width=8, anchor="w").pack(side="left")
        self.plot_min_freq_entry = tk.Entry(freq_row, textvariable=self.plot_min_freq_var, width=7)
        self.plot_min_freq_entry.pack(side="left", padx=2)
        tk.Label(freq_row, text="to").pack(side="left", padx=2)
        self.plot_max_freq_entry = tk.Entry(freq_row, textvariable=self.plot_max_freq_var, width=7)
        self.plot_max_freq_entry.pack(side="left", padx=2)
        tk.Label(freq_row, text="Hz", fg="gray").pack(side="left", padx=2)
        
        # Initialize entry states based on auto-limit
        self._toggle_vel_limits()
        self._toggle_freq_limits()
        
        # Tick spacing options
        tick_label = tk.Label(image_frame, text="Axis Tick Spacing:", anchor="w")
        tick_label.pack(fill="x", pady=(6, 2))
        
        freq_tick_row = tk.Frame(image_frame); freq_tick_row.pack(fill="x", pady=2)
        tk.Label(freq_tick_row, text="Frequency:", width=12, anchor="w").pack(side="left")
        ttk.Combobox(freq_tick_row, textvariable=self.freq_tick_spacing_var,
                     values=["auto", "1", "2", "5", "10", "20"],
                     width=8).pack(side="left", padx=2)
        tk.Label(freq_tick_row, text="Hz", fg="gray").pack(side="left", padx=2)
        
        vel_tick_row = tk.Frame(image_frame); vel_tick_row.pack(fill="x", pady=2)
        tk.Label(vel_tick_row, text="Velocity:", width=12, anchor="w").pack(side="left")
        ttk.Combobox(vel_tick_row, textvariable=self.vel_tick_spacing_var,
                     values=["auto", "10", "25", "50", "100", "200"],
                     width=8).pack(side="left", padx=2)
        tk.Label(vel_tick_row, text="m/s", fg="gray").pack(side="left", padx=2)
        
        # Colormap
        cmap_row = tk.Frame(image_frame); cmap_row.pack(fill="x", pady=3)
        tk.Label(cmap_row, text="Colormap:", width=12, anchor="w").pack(side="left")
        ttk.Combobox(cmap_row, textvariable=self.cmap_var,
                     values=["jet", "viridis", "plasma", "turbo", "seismic", "hot"],
                     width=10, state="readonly").pack(side="left", padx=4)
        
        # DPI
        dpi_row = tk.Frame(image_frame); dpi_row.pack(fill="x", pady=3)
        tk.Label(dpi_row, text="DPI:", width=16, anchor="w").pack(side="left")
        tk.Entry(dpi_row, textvariable=self.dpi_var, width=8).pack(side="left", padx=4)
        tk.Label(dpi_row, text="(72-600)", fg="gray").pack(side="left")
        
        # Export spectra checkbox
        exp_row = tk.Frame(image_frame); exp_row.pack(fill="x", pady=3)
        tk.Checkbutton(exp_row, text="Export power spectra (.npz)", 
                       variable=self.export_spectra_var).pack(side="left")
        
        # ===== Buttons =====
        btn_frame = tk.Frame(scrollable)
        btn_frame.pack(fill="x", padx=8, pady=12)
        
        tk.Button(btn_frame, text="Reset to Defaults", 
                  command=self._reset_advanced_defaults).pack(side="left", padx=4)
        tk.Button(btn_frame, text="Close", 
                  command=popup.destroy).pack(side="right", padx=4)
        
        # Center the window
        popup.update_idletasks()
        x = self.root.winfo_rootx() + (self.root.winfo_width() - popup.winfo_width()) // 2
        y = self.root.winfo_rooty() + (self.root.winfo_height() - popup.winfo_height()) // 2
        popup.geometry(f"+{x}+{y}")

    def _edit_cell(self, ev):
        """Handle double-click editing in file tree.
        
        Columns: (file, type, offset, rev)
        - Column 1 (file): not editable
        - Column 2 (type): not editable
        - Column 3 (offset): editable text entry
        - Column 4 (rev): toggle checkbox
        """
        item = self.tree.identify_row(ev.y); col = self.tree.identify_column(ev.x)
        if not item: return
        col_idx = int(col.lstrip("#")); vals = list(self.tree.item(item, "values")); base = vals[0]
        x, y, w, h = self.tree.bbox(item, col)
        if col_idx == 3:  # Offset column
            e = tk.Entry(self.root); e.insert(0, vals[2])
            e.place(x=x, y=y + self.tree.winfo_rooty()-self.root.winfo_rooty(), width=w, height=h); e.focus()
            def _done(_):
                vals[2] = e.get(); self.offsets[base] = vals[2]; self.tree.item(item, values=vals); e.destroy()
            e.bind("<Return>", _done); e.bind("<FocusOut>", _done)
        elif col_idx == 4:  # Reverse column
            flag = vals[3] != "☑"; self.reverse_flags[base] = flag; vals[3] = "☑" if flag else "☐"; self.tree.item(item, values=vals)

    def _collect_common(self):
        try:
            pick_vmin = float(self.vmin_var.get()); pick_vmax = float(self.vmax_var.get())
            pick_fmin = float(self.fmin_var.get()); pick_fmax = float(self.fmax_var.get())
        except ValueError:
            raise ValueError("Invalid pick/plot limits")
        try:
            st = float(self.time_start_var.get()); en = float(self.time_end_var.get())
            downsample = bool(self.downsample_var.get()); dfac = int(self.down_factor_var.get()); numf = int(self.numf_var.get())
        except Exception as e:
            raise ValueError(f"Invalid sampling: {e}")
        try:
            fig_dpi = int(self.dpi_var.get()); assert fig_dpi > 0
        except Exception:
            raise ValueError("Figure DPI must be a positive integer")
        return pick_vmin, pick_vmax, pick_fmin, pick_fmax, st, en, downsample, dfac, numf, fig_dpi

    def _resolve_selected_path(self) -> str | None:
        if not self.file_list:
            return None
        sel = self.tree.selection()
        if sel:
            want = self.tree.item(sel[0], "values")[0]
            for f in self.file_list:
                if os.path.splitext(os.path.basename(f))[0] == want:
                    return f
        return self.file_list[0]

    def preview_array(self):
        # Build an embedded preview (array schematic + waterfall)
        path = self._resolve_selected_path()
        if path is None:
            messagebox.showerror("Preview", "No file selected."); return
        try:
            from matplotlib.figure import Figure
            from matplotlib.backends.backend_tkagg import FigureCanvasTkAgg
            from sw_transform.processing.seg2 import load_seg2_ar
            from sw_transform.processing.preprocess import preprocess_data
            import numpy as np
            time, T, Shotpoint, Spacing, dt, _ = load_seg2_ar(path)
            st = float(self.time_start_var.get()); en = float(self.time_end_var.get())
            ds = bool(self.downsample_var.get()); df = int(self.down_factor_var.get()); nf = int(self.numf_var.get())
            Tpre, time_pre, dt2 = preprocess_data(T, time, dt, reverse_shot=False, start_time=st, end_time=en, do_downsample=ds, down_factor=df, numf=nf)
            # Optional display time trim
            disp_txt = (self.display_time_var.get() or "").strip()
            if disp_txt:
                try:
                    disp_time = float(disp_txt)
                    if disp_time > 0:
                        n_keep = int(np.clip(disp_time / dt2, 1, Tpre.shape[0]))
                        Tpre = Tpre[:n_keep, :]; time_pre = time_pre[:n_keep]
                except Exception:
                    pass
            positions = np.arange(Tpre.shape[1], dtype=float) * float(Spacing)
            fig = Figure(figsize=(7.5, 6.0), dpi=100)
            gs = fig.add_gridspec(2, 1, height_ratios=[1, 3], hspace=0.42)
            ax1 = fig.add_subplot(gs[0]); ax2 = fig.add_subplot(gs[1])
            # schematic
            ax1.plot(positions, np.zeros_like(positions), "^", color="green", label="Sensor")
            # Source position: prefer table override
            try:
                base = os.path.splitext(os.path.basename(path))[0]
                off_txt = (self.offsets.get(base, "+0") or "+0").strip().replace("m", "")
                if off_txt.startswith("+"):
                    src_x = float(off_txt[1:])
                else:
                    src_x = float(off_txt)
            except Exception:
                src_x = float(Shotpoint)
            ax1.plot([src_x], [0.0], "D", color="tab:red", label="Source")
            ax1.set_yticks([]); ax1.set_xlabel("Distance (m)")
            ax1.legend(loc="upper left", bbox_to_anchor=(1.02, 1), borderaxespad=0)
            ax1.set_title("Array schematic")
            # waterfall
            traces = Tpre.copy().T
            denom = np.max(np.abs(traces), axis=1, keepdims=True); denom[denom==0]=1.0; traces = traces/denom
            spacing = float(np.mean(np.diff(positions))) if len(positions)>1 else 1.0
            scale = 0.5 * spacing
            for tr, x0 in zip(traces, positions):
                ax2.plot(tr*scale + x0, time_pre, color="b", linewidth=0.5)
            ax2.invert_yaxis(); ax2.set_xlabel("Distance (m)"); ax2.set_ylabel("Time (s)"); ax2.set_title("Waterfall (normalized)")
            fig.tight_layout(rect=[0, 0, 0.88, 1])  # Leave space for legend on right
            # mount
            if self.prev_canvas_widget is not None:
                try: self.prev_canvas_widget.destroy()
                except Exception: pass
            canvas = FigureCanvasTkAgg(fig, master=self.prev_host)
            self.prev_canvas_widget = canvas.get_tk_widget(); self.prev_canvas_widget.pack(fill="both", expand=True)
            canvas.draw()
        except Exception as e:
            messagebox.showerror("Preview", str(e))

    def _get_worker_count(self, mode: str = 'single') -> int:
        """Get worker count based on user setting or auto mode."""
        from sw_transform.workers.parallel import get_optimal_workers
        val = self.worker_count_var.get()
        if val == "auto":
            return get_optimal_workers(mode=mode)
        try:
            return max(1, int(val))
        except ValueError:
            return get_optimal_workers(mode=mode)

    def run_single_processing(self, selected_only: bool=False):
        if not self.file_list:
            messagebox.showerror("No files", "Select data files first."); return
        if not self.output_folder:
            messagebox.showerror("No output", "Choose an output folder."); return
        key = self.method_key.get()
        
        # Check if any .mat files and method is not FDBF
        has_mat = any(self.file_types.get(os.path.splitext(os.path.basename(p))[0]) == 'mat' 
                      for p in self.file_list)
        if has_mat and key != 'fdbf':
            messagebox.showerror("Method Error", 
                "Vibrosis .mat files only support FDBF method.\n"
                "Please select FDBF or remove .mat files.")
            return
        
        if key in ("fk","fdbf"):
            grid_n = int(self.grid_fk_var.get()); tol = float(self.tol_fk_var.get()); vspace = "linear"
        else:
            grid_n = int(self.grid_ps_var.get()); tol = float(self.tol_ps_var.get()); vspace = self.vspace_ps_var.get() or "linear"
        pick_vmin, pick_vmax, pick_fmin, pick_fmax, st, en, downsample, dfac, numf, fig_dpi = self._collect_common()
        
        # Get dx for vibrosis .mat files
        try:
            dx = float(self.dx_var.get())
        except ValueError:
            dx = 1.0
        
        try:
            import matplotlib as mpl; old_dpi = mpl.rcParams.get('savefig.dpi', 'figure'); mpl.rcParams['savefig.dpi'] = fig_dpi
        except Exception:
            old_dpi = None
        try:
            paths = list(self.file_list)
            if selected_only:
                sel = self.tree.selection(); selected_bases = {self.tree.item(i, "values")[0] for i in sel} if sel else set()
                paths = [p for p in self.file_list if os.path.splitext(os.path.basename(p))[0] in selected_bases]
            
            # Build params list for all files
            params_list = []
            for path in paths:
                base = os.path.splitext(os.path.basename(path))[0]
                offset = self.offsets.get(base, "+0")
                user_rev = self.reverse_flags.get(base, False)
                rev = compute_reverse_flag(bool(user_rev), key)
                source_type = "vibrosis" if self.vibrosis_mode.get() else "hammer"
                cylindrical = self.cylindrical_var.get()
                file_type = self.file_types.get(base, 'seg2')
                
                params = dict(path=path, base=base, key=key, offset=offset, outdir=self.output_folder,
                              pick_vmin=pick_vmin, pick_vmax=pick_vmax, pick_fmin=pick_fmin, pick_fmax=pick_fmax,
                              st=st, en=en, downsample=downsample, dfac=dfac, numf=numf,
                              grid_n=grid_n, tol=tol, vspace=vspace, dpi=fig_dpi, rev=rev,
                              topic=(self.figure_topic_var.get() or ""),
                              source_type=source_type,
                              cylindrical=cylindrical,
                              export_spectra=self.export_spectra_var.get(),
                              file_type=file_type,
                              dx=dx if file_type == 'mat' else None,
                              auto_vel_limits=self.auto_vel_limits_var.get(),
                              auto_freq_limits=self.auto_freq_limits_var.get(),
                              plot_min_vel=self.plot_min_vel_var.get(),
                              plot_max_vel=self.plot_max_vel_var.get(),
                              plot_min_freq=self.plot_min_freq_var.get(),
                              plot_max_freq=self.plot_max_freq_var.get(),
                              cmap=self.cmap_var.get(),
                              freq_tick_spacing=self.freq_tick_spacing_var.get(),
                              vel_tick_spacing=self.vel_tick_spacing_var.get())
                params_list.append(params)
            
            total = max(1, len(params_list))
            self.pb.config(maximum=total, value=0)
            use_parallel = self.parallel_var.get() and len(params_list) > 1
            
            if use_parallel:
                # Parallel processing
                from sw_transform.workers.parallel import run_batch_parallel
                n_workers = self._get_worker_count(mode='single')
                self.pb_label.config(text=f"Processing {total} files ({n_workers} workers)...")
                self.root.update()
                if self.logbox:
                    self.logbox.insert(tk.END, f"Starting parallel processing: {total} files on {n_workers} workers\n")
                    self.logbox.see(tk.END)
                
                # Progress callback updates GUI
                def progress_cb(completed, total_count, current_base):
                    self.pb['value'] = completed
                    self.pb_label.config(text=f"Completed {completed}/{total_count}: {current_base}")
                    if self.logbox:
                        self.logbox.insert(tk.END, f"Finished: {current_base}\n")
                        self.logbox.see(tk.END)
                    self.root.update()
                
                results = run_batch_parallel(params_list, mode='single', max_workers=n_workers, progress_callback=progress_cb)
                success = sum(1 for r in results if r.success)
                failures = sum(1 for r in results if not r.success)
                
                # Log any failures
                for r in results:
                    if not r.success and self.logbox:
                        self.logbox.insert(tk.END, f"Failed {r.base}: {r.error}\n")
                        self.logbox.see(tk.END)
            else:
                # Sequential processing (original code)
                self.pb_label.config(text="Processing…")
                self.root.update_idletasks()
                success = 0; failures = 0
                for idx, params in enumerate(params_list):
                    base = params['base']
                    self.pb_label.config(text=f"Processing {idx+1}/{total}: {base}")
                    self.pb['value'] = idx
                    self.root.update()
                    if self.logbox: self.logbox.insert(tk.END, f"Running {base} [{key}]...\n"); self.logbox.see(tk.END)
                    _b, ok, out = svc_run_single(params)
                    if ok and isinstance(out, str) and out.lower().endswith('.png'):
                        success += 1
                        if self.logbox: self.logbox.insert(tk.END, f"Saved: {out}\n"); self.logbox.see(tk.END)
                    else:
                        failures += 1
                        if self.logbox: self.logbox.insert(tk.END, f"Failed {base}: {out}\n"); self.logbox.see(tk.END)
                    self.pb['value'] = idx + 1
                    self.root.update()
            
            # After all files processed, create combined CSV
            if success > 0:
                try:
                    self._create_combined_csv_single_method(key, paths)
                    if self.logbox:
                        self.logbox.insert(tk.END, f"Created combined_{key}.csv\n")
                        self.logbox.see(tk.END)
                except Exception as e:
                    if self.logbox:
                        self.logbox.insert(tk.END, f"Warning: Could not create combined CSV: {e}\n")
                        self.logbox.see(tk.END)
                # Create combined spectrum if export was enabled
                if self.export_spectra_var.get() and success > 1:
                    try:
                        self._create_combined_spectrum_single_method(key, paths)
                        if self.logbox:
                            self.logbox.insert(tk.END, f"Created combined_{key}_spectrum.npz\n")
                            self.logbox.see(tk.END)
                    except Exception as e:
                        if self.logbox:
                            self.logbox.insert(tk.END, f"Warning: Could not create combined spectrum: {e}\n")
                            self.logbox.see(tk.END)
            if success > 0 and failures == 0:
                messagebox.showinfo("Run", f"Completed ({success} file(s))")
            elif success > 0:
                messagebox.showwarning("Run", f"Completed with errors: {success} success, {failures} failed. See log.")
            else:
                messagebox.showerror("Run", "No outputs created. Check the log for errors.")
            # Optional PPT
            if success > 0 and bool(self.ppt_var.get()):
                self._build_ppt_from_gallery()
            # Build combined CSV across all shot offsets
            if success > 1:
                self._build_combined_csv_single(key, paths)
        except Exception as e:
            messagebox.showerror("Run", str(e))
        finally:
            try:
                import matplotlib as mpl
                if old_dpi is not None:
                    mpl.rcParams['savefig.dpi'] = old_dpi
            except Exception:
                pass
            self.pb_label.config(text="Idle"); self.pb.config(value=0)

    def run_compare_processing(self, selected_only: bool=False):
        if not self.file_list:
            messagebox.showerror("No files", "Select data files first."); return
        if not self.output_folder:
            messagebox.showerror("No output", "Choose an output folder."); return
        
        # Check if any .mat files - compare mode not supported for .mat
        has_mat = any(self.file_types.get(os.path.splitext(os.path.basename(p))[0]) == 'mat' 
                      for p in self.file_list)
        if has_mat:
            messagebox.showerror("Compare Not Supported", 
                "Compare mode is not supported for vibrosis .mat files.\n"
                "Use 'Run Selected/All' with FDBF method instead.")
            return
        
        pick_vmin, pick_vmax, pick_fmin, pick_fmax, st, en, downsample, dfac, numf, fig_dpi = self._collect_common()
        n_fk = int(self.grid_fk_var.get()); tol_fk = float(self.tol_fk_var.get())
        n_ps = int(self.grid_ps_var.get()); vspace_ps = self.vspace_ps_var.get() or "linear"
        paths = list(self.file_list)
        if selected_only:
            sel = self.tree.selection(); selected_bases = {self.tree.item(i, "values")[0] for i in sel} if sel else set()
            paths = [p for p in self.file_list if os.path.splitext(os.path.basename(p))[0] in selected_bases]
        
        # Build params list for all files
        params_list = []
        for path in paths:
            base = os.path.splitext(os.path.basename(path))[0]
            offset = self.offsets.get(base, "+0")
            user_rev = self.reverse_flags.get(base, False)
            source_type = "vibrosis" if self.vibrosis_mode.get() else "hammer"
            cylindrical = self.cylindrical_var.get()
            params = dict(path=path, base=base, outdir=self.output_folder, offset=offset,
                          pick_vmin=pick_vmin, pick_vmax=pick_vmax, pick_fmin=pick_fmin, pick_fmax=pick_fmax,
                          st=st, en=en, downsample=downsample, dfac=dfac, numf=numf,
                          n_fk=n_fk, tol_fk=tol_fk, n_ps=n_ps, vspace_ps=vspace_ps,
                          rev_fk=user_rev, rev_ps=user_rev, rev_fdbf=user_rev, rev_ss=user_rev,
                          topic=(self.figure_topic_var.get() or ""),
                          source_type=source_type,
                          cylindrical=cylindrical,
                          export_spectra=self.export_spectra_var.get())
            params_list.append(params)
        
        total = max(1, len(params_list))
        self.pb.config(maximum=total, value=0)
        use_parallel = self.parallel_var.get() and len(params_list) > 1
        
        if use_parallel:
            # Parallel processing - use fewer workers for compare (memory intensive)
            from sw_transform.workers.parallel import run_batch_parallel
            n_workers = self._get_worker_count(mode='compare')
            self.pb_label.config(text=f"Comparing {total} files ({n_workers} workers)...")
            self.root.update()
            if self.logbox:
                self.logbox.insert(tk.END, f"Starting parallel comparison: {total} files on {n_workers} workers\n")
                self.logbox.see(tk.END)
            
            def progress_cb(completed, total_count, current_base):
                self.pb['value'] = completed
                self.pb_label.config(text=f"Completed {completed}/{total_count}: {current_base}")
                if self.logbox:
                    self.logbox.insert(tk.END, f"Finished: {current_base}\n")
                    self.logbox.see(tk.END)
                self.root.update()
            
            results = run_batch_parallel(params_list, mode='compare', max_workers=n_workers, progress_callback=progress_cb)
            success = sum(1 for r in results if r.success)
            failures = sum(1 for r in results if not r.success)
            
            for r in results:
                if not r.success and self.logbox:
                    self.logbox.insert(tk.END, f"Failed {r.base}: {r.error}\n")
                    self.logbox.see(tk.END)
        else:
            # Sequential processing
            self.pb_label.config(text="Comparing…")
            self.root.update_idletasks()
            success = 0; failures = 0
            for idx, params in enumerate(params_list):
                base = params['base']
                self.pb_label.config(text=f"Comparing {idx+1}/{total}: {base}")
                self.pb['value'] = idx
                self.root.update()
                if self.logbox: self.logbox.insert(tk.END, f"Compare {base}...\n"); self.logbox.see(tk.END)
                _b, ok, out = svc_run_compare(params)
                if ok and isinstance(out, str) and out.lower().endswith('.png'):
                    success += 1
                    if self.logbox: self.logbox.insert(tk.END, f"Saved: {out}\n"); self.logbox.see(tk.END)
                else:
                    failures += 1
                    if self.logbox: self.logbox.insert(tk.END, f"Failed {base}: {out}\n"); self.logbox.see(tk.END)
                self.pb['value'] = idx + 1
                self.root.update()
        # After all files processed, create combined CSV for compare mode
        if success > 0:
            try:
                self._create_combined_csv_compare_mode(paths)
                if self.logbox:
                    self.logbox.insert(tk.END, "Created combined_compare_all.csv\n")
                    self.logbox.see(tk.END)
            except Exception as e:
                if self.logbox:
                    self.logbox.insert(tk.END, f"Warning: Could not create combined CSV: {e}\n")
                    self.logbox.see(tk.END)
            # Create combined spectra for all 4 methods if export was enabled
            if self.export_spectra_var.get() and success > 1:
                for method in ['fk', 'fdbf', 'ps', 'ss']:
                    try:
                        self._create_combined_spectrum_single_method(method, paths)
                        if self.logbox:
                            self.logbox.insert(tk.END, f"Created combined_{method}_spectrum.npz\n")
                            self.logbox.see(tk.END)
                    except Exception as e:
                        if self.logbox:
                            self.logbox.insert(tk.END, f"Warning: Could not create combined {method} spectrum: {e}\n")
                            self.logbox.see(tk.END)
        if success > 0 and failures == 0:
            messagebox.showinfo("Run", f"Comparison completed ({success} file(s))")
        elif success > 0:
            messagebox.showwarning("Run", f"Comparison completed with errors: {success} success, {failures} failed. See log.")
        else:
            messagebox.showerror("Run", "No comparison outputs created. Check the log for errors.")
        # Optional PPT
        if success > 0 and bool(self.ppt_var.get()):
            self._build_ppt_from_gallery()
        # Build global combined CSV across all files (4 methods)
        if success > 1:
            self._build_combined_csv_compare(paths)
        self.pb_label.config(text="Idle"); self.pb.config(value=0)

    def _create_combined_csv_single_method(self, key: str, paths: list):
        """Aggregate all per-shot CSVs into combined_<method>.csv."""
        import csv
        import glob
        # Find all per-shot CSVs for this method
        pattern = os.path.join(self.output_folder, f"*_{key}_*.csv")
        csv_files = glob.glob(pattern)
        # Exclude any compare or combined files
        csv_files = [f for f in csv_files if not ('_compare' in f or 'combined_' in os.path.basename(f))]
        if not csv_files:
            return
        # Read all CSVs and organize by (base, offset)
        data_map = {}  # (base, offset) -> (freq, vel, wav)
        for csv_path in csv_files:
            try:
                with open(csv_path, 'r', newline='') as f:
                    reader = csv.reader(f)
                    header = next(reader, None)
                    if not header or len(header) < 2:
                        continue
                    freq_list = []; vel_list = []; wav_list = []
                    for row in reader:
                        if len(row) >= 2:
                            freq_list.append(row[0])
                            vel_list.append(row[1])
                            wav_list.append(row[2] if len(row) > 2 else "")
                    # Extract base and offset from filename
                    fname = os.path.basename(csv_path)
                    # Format: <base>_<method>_<offset_tag>.csv
                    # e.g., "file1_fk_p66.csv" -> base="file1", offset="+66"
                    parts = fname.replace(".csv", "").split("_")
                    if len(parts) >= 3:
                        # Find the method position in the filename
                        method_idx = -1
                        for idx, part in enumerate(parts):
                            if part == key:
                                method_idx = idx
                                break
                        if method_idx >= 0 and method_idx < len(parts) - 1:
                            base = "_".join(parts[:method_idx])
                            offset_tag = parts[method_idx + 1]
                            # Convert offset_tag back to display format
                            if offset_tag.startswith('p'):
                                offset = "+" + offset_tag[1:] + "m"
                            elif offset_tag.startswith('m'):
                                offset = "-" + offset_tag[1:] + "m"
                            else:
                                offset = offset_tag
                            data_map[(base, offset)] = (freq_list, vel_list, wav_list)
            except Exception as e:
                if self.logbox:
                    self.logbox.insert(tk.END, f"Warning: Failed to read {csv_path}: {e}\n")
                continue
        if not data_map:
            return
        # Create combined CSV
        combined_path = os.path.join(self.output_folder, f"combined_{key}.csv")
        with open(combined_path, 'w', newline='') as f:
            writer = csv.writer(f)
            # Build header
            header = []
            sorted_keys = sorted(data_map.keys())
            for base, offset in sorted_keys:
                header.extend([f"freq({key}_{offset})", f"vel({key}_{offset})", f"wav({key}_{offset})"])
            writer.writerow(header)
            # Find max length
            max_len = max(len(data_map[k][0]) for k in sorted_keys) if sorted_keys else 0
            # Write rows
            for i in range(max_len):
                row = []
                for base, offset in sorted_keys:
                    freq_list, vel_list, wav_list = data_map[(base, offset)]
                    row.append(freq_list[i] if i < len(freq_list) else "")
                    row.append(vel_list[i] if i < len(vel_list) else "")
                    row.append(wav_list[i] if i < len(wav_list) else "")
                writer.writerow(row)

    def _create_combined_csv_compare_mode(self, paths: list):
        """Aggregate all per-file compare CSVs into combined_compare_all.csv."""
        import csv
        import glob
        # Find all per-file compare CSVs
        pattern = os.path.join(self.output_folder, "*_compare.csv")
        csv_files = glob.glob(pattern)
        if not csv_files:
            return
        # Read all CSVs
        data_map = {}  # base -> list of rows
        headers_map = {}  # base -> header
        for csv_path in csv_files:
            try:
                with open(csv_path, 'r', newline='') as f:
                    reader = csv.reader(f)
                    header = next(reader, None)
                    if not header:
                        continue
                    rows = list(reader)
                    base = os.path.basename(csv_path).replace("_compare.csv", "")
                    data_map[base] = rows
                    headers_map[base] = header
            except Exception:
                continue
        if not data_map:
            return
        # Create combined CSV
        combined_path = os.path.join(self.output_folder, "combined_compare_all.csv")
        with open(combined_path, 'w', newline='') as f:
            writer = csv.writer(f)
            # Build combined header
            sorted_bases = sorted(data_map.keys())
            combined_header = []
            for base in sorted_bases:
                combined_header.extend(headers_map.get(base, []))
            writer.writerow(combined_header)
            # Find max length
            max_len = max(len(data_map[b]) for b in sorted_bases) if sorted_bases else 0
            # Write rows
            for i in range(max_len):
                row = []
                for base in sorted_bases:
                    file_rows = data_map[base]
                    if i < len(file_rows):
                        row.extend(file_rows[i])
                    else:
                        # Pad with empty cells matching header width
                        row.extend([""] * len(headers_map.get(base, [])))
                writer.writerow(row)

    def _create_combined_spectrum_single_method(self, key: str, paths: list):
        """Create combined spectrum .npz file for a single method across all source offsets."""
        import glob
        from sw_transform.core.service import create_combined_spectrum

        # Find all individual spectrum files for this method
        pattern = os.path.join(self.output_folder, f"*_{key}_*_spectrum.npz")
        spectrum_files = glob.glob(pattern)

        # Exclude combined files
        spectrum_files = [f for f in spectrum_files if not os.path.basename(f).startswith('combined_')]

        if not spectrum_files:
            return

        # Create combined spectrum
        create_combined_spectrum(self.output_folder, key, spectrum_files)

    # data helpers
    def _auto_assign_from_filenames(self):
        if not self.file_list:
            messagebox.showerror("No files", "Select SEG-2 files first."); return
        try:
            rows = assign_files_from_names(self.file_list, recursive=False, include_unknown=True)
            new_offsets = {}; new_reverse = {}
            for r in rows:
                try:
                    fp = str(getattr(r, 'file_path', ''))
                    b = os.path.splitext(os.path.basename(fp))[0]
                    off = getattr(r, 'offset_m', None); rev = bool(getattr(r, 'reverse', False))
                    if off is not None:
                        val = float(off); sign = "+" if val >= 0 else ""; new_offsets[b] = f"{sign}{int(round(val))}"
                    new_reverse[b] = rev
                except Exception:
                    continue
            for item in self.tree.get_children():
                vals = list(self.tree.item(item, "values")); base = vals[0]
                if base in new_offsets:
                    vals[1] = new_offsets[base]; self.offsets[base] = vals[1]
                if base in new_reverse:
                    self.reverse_flags[base] = bool(new_reverse[base]); vals[2] = "☑" if self.reverse_flags[base] else "☐"
                self.tree.item(item, values=vals)
            messagebox.showinfo("Assign", "Offsets/reverse assigned from filenames.")
        except Exception as e:
            messagebox.showerror("Assign", str(e))

    # --- Figures tab helpers ---
    def refresh_gallery(self):
        if not self.output_folder:
            messagebox.showwarning("Figures", "Select an output folder first."); return
        try:
            import glob, os
            self.fig_list.delete(*self.fig_list.get_children())
            for path in glob.glob(os.path.join(self.output_folder, "**", "*.png"), recursive=True):
                name = os.path.basename(path)
                self.fig_list.insert("", "end", values=(name, path))
        except Exception as e:
            messagebox.showerror("Figures", str(e))

    def _on_figure_selected(self, _):
        sel = self.fig_list.selection()
        if not sel:
            return
        path = self.fig_list.item(sel[0], "values")[1]
        self._load_preview_image(path)
        self._render_fig_preview()

    def _load_preview_image(self, path: str):
        try:
            from PIL import Image, ImageTk
        except Exception:
            self._fig_image_pil = None; self._fig_image_tk = None; return
        try:
            self._fig_image_pil = Image.open(path)
            self._fig_scale = 1.0
            self.fig_fit_mode.set("Auto")
        except Exception:
            self._fig_image_pil = None; self._fig_image_tk = None

    def _render_fig_preview(self):
        if self._fig_image_pil is None:
            self.fig_prev_canvas.delete("all"); self.fig_prev_canvas.config(scrollregion=(0,0,0,0)); return
        try:
            from PIL import ImageTk
            cw = self.fig_prev_canvas.winfo_width(); ch = self.fig_prev_canvas.winfo_height()
            if cw <= 1 or ch <= 1:
                self.root.update_idletasks(); cw = self.fig_prev_canvas.winfo_width() or 900; ch = self.fig_prev_canvas.winfo_height() or 600
            ow, oh = self._fig_image_pil.size
            mode = (self.fig_fit_mode.get() or "Auto").strip()
            if mode == "Width":
                scale = max(1e-3, (cw-8)/ow)
            elif mode == "Height":
                scale = max(1e-3, (ch-8)/oh)
            elif mode == "Auto":
                scale = max(1e-3, min((cw-8)/ow, (ch-8)/oh))
            else:
                scale = max(self._fig_scale, 1e-3)
            tw = max(50, int(ow*scale)); th = max(50, int(oh*scale))
            self._fig_image_tk = ImageTk.PhotoImage(self._fig_image_pil.resize((tw, th)))
            self.fig_prev_canvas.delete("all")
            self.fig_prev_canvas.create_image(0, 0, image=self._fig_image_tk, anchor="nw")
            self.fig_prev_canvas.config(scrollregion=(0,0,tw,th))
            try:
                self.fig_zoom_label.set(f"{int(round(scale*100))}%")
            except Exception:
                pass
        except Exception:
            pass

    def _fig_zoom_step(self, factor: float):
        self.fig_fit_mode.set("None")
        self._fig_scale *= float(factor)
        if self._fig_scale < 0.05:
            self._fig_scale = 0.05
        self._render_fig_preview()

    def _fig_zoom_reset(self):
        self._fig_scale = 1.0
        self.fig_fit_mode.set("Auto")
        self._render_fig_preview()

    def _on_fig_wheel_zoom(self, event):
        step = 1.05
        self.fig_fit_mode.set("None")
        if event.delta > 0:
            self._fig_scale *= step
        else:
            self._fig_scale /= step
            if self._fig_scale < 0.05:
                self._fig_scale = 0.05
        self._render_fig_preview()

    # --- icon helper ---
    def _load_icon(self, name: str, size: int) -> tk.PhotoImage | None:
        """Load and cache an icon using the utility function."""
        return load_icon(self._icons, name, size)

    def _open_selected_figure(self):
        sel = self.fig_list.selection()
        if not sel: return
        path = self.fig_list.item(sel[0], "values")[1]
        try:
            import subprocess, sys
            if os.name == 'nt':
                os.startfile(path)
            elif sys.platform == 'darwin':
                subprocess.call(['open', path])
            else:
                subprocess.call(['xdg-open', path])
        except Exception:
            pass

    def _delete_selected_figure(self):
        sel = self.fig_list.selection()
        if not sel: return
        path = self.fig_list.item(sel[0], "values")[1]
        try:
            os.remove(path)
            self.refresh_gallery()
        except Exception:
            pass

    def _build_ppt_from_gallery(self):
        if not self.output_folder:
            messagebox.showwarning("PowerPoint", "Select an output folder first."); return
        try:
            from pptx import Presentation
            from pptx.util import Inches
        except Exception:
            messagebox.showwarning("PowerPoint", "python-pptx not installed. Run: pip install python-pptx"); return
        try:
            import glob, os
            pngs = glob.glob(os.path.join(self.output_folder, "**", "*.png"), recursive=True)
            if not pngs:
                messagebox.showinfo("PowerPoint", "No PNG files found in output folder."); return
            out_ppt = os.path.join(self.output_folder, "slides_all_outputs.pptx")
            prs = Presentation()
            blank = prs.slide_layouts[6]
            sw_in = float(prs.slide_width) / 914400.0
            sh_in = float(prs.slide_height) / 914400.0
            tw_in = max(0.5, sw_in - 1.0)
            th_in = max(0.5, sh_in - 1.0)
            for img in pngs:
                slide = prs.slides.add_slide(blank)
                try:
                    from PIL import Image
                    Image.MAX_IMAGE_PIXELS = None
                    with Image.open(img) as im: iw, ih = im.size
                    ar = iw/ih; target_ar = tw_in/max(th_in, 1e-6)
                    if target_ar > ar:
                        h_in = th_in; w_in = h_in*ar
                    else:
                        w_in = tw_in; h_in = w_in/max(ar, 1e-9)
                    x = Inches(max(0.0, (sw_in - w_in) / 2.0)); y = Inches(max(0.0, (sh_in - h_in) / 2.0))
                    slide.shapes.add_picture(img, x, y, width=Inches(w_in), height=Inches(h_in))
                except Exception:
                    pass
            prs.save(out_ppt)
            messagebox.showinfo("PowerPoint", f"Created: {out_ppt}")
            self.refresh_gallery()
        except Exception as e:
            messagebox.showerror("PowerPoint", str(e))

    def _build_combined_csv_single(self, key: str, paths: list[str]):
        """Aggregate per-shot CSVs into combined_<method>.csv across all offsets."""
        try:
            import csv
            combined_rows = []
            for path in paths:
                base = os.path.splitext(os.path.basename(path))[0]
                offset = self.offsets.get(base, "+0")
                _off = str(offset).strip().replace(" ", "").replace("m", "")
                if _off.startswith("+"):
                    _off_tag = "p" + _off[1:]
                elif _off.startswith("-"):
                    _off_tag = "m" + _off[1:]
                else:
                    _off_tag = _off
                per_csv = os.path.join(self.output_folder, f"{base}_{key}_{_off_tag}.csv")
                if not os.path.isfile(per_csv):
                    continue
                frq, vel, wl = [], [], []
                try:
                    with open(per_csv, "r", newline="") as fcsv:
                        rdr = csv.reader(fcsv)
                        next(rdr, None)
                        for row in rdr:
                            if len(row) >= 3:
                                try:
                                    frq.append(float(row[0])); vel.append(float(row[1])); wl.append(float(row[2]))
                                except Exception:
                                    continue
                except Exception:
                    continue
                if frq:
                    combined_rows.append((offset, frq, vel, wl))
            if combined_rows and len(combined_rows) > 1:
                min_len = min(len(r[1]) for r in combined_rows)
                comb_csv = os.path.join(self.output_folder, f"combined_{key}.csv")
                with open(comb_csv, "w", newline="") as fcsv:
                    w = csv.writer(fcsv)
                    header = []
                    for off, _, _, _ in combined_rows:
                        header += [f"freq({key}_{off})", f"vel({key}_{off})", f"wav({key}_{off})"]
                    w.writerow(header)
                    for i in range(min_len):
                        row = []
                        for _, frq, vel, wl in combined_rows:
                            row += [frq[i], vel[i], wl[i] if i < len(wl) else ""]
                        w.writerow(row)
                if self.logbox:
                    self.logbox.insert(tk.END, f"Combined CSV → {comb_csv}\n"); self.logbox.see(tk.END)
        except Exception:
            pass

    def _build_combined_csv_compare(self, paths: list[str]):
        """Aggregate per-file compare CSVs into combined_compare_all.csv."""
        try:
            import csv
            rows_by_file = []
            headers = []
            for pth in paths:
                b = os.path.splitext(os.path.basename(pth))[0]
                comb_path = os.path.join(self.output_folder, f"{b}_compare.csv")
                if not os.path.isfile(comb_path):
                    continue
                with open(comb_path, "r", newline="") as fcsv:
                    rdr = list(csv.reader(fcsv))
                    if not rdr:
                        continue
                    hdr = rdr[0]; data = rdr[1:]
                    headers += [f"{h}[{b}]" for h in hdr]
                    rows_by_file.append(data)
            if rows_by_file and len(rows_by_file) > 1:
                min_len = min(len(r) for r in rows_by_file)
                out_all = os.path.join(self.output_folder, "combined_compare_all.csv")
                with open(out_all, "w", newline="") as fcsv:
                    w = csv.writer(fcsv)
                    w.writerow(headers)
                    for i in range(min_len):
                        row = []
                        for data in rows_by_file:
                            row += data[i]
                        w.writerow(row)
                if self.logbox:
                    self.logbox.insert(tk.END, f"Combined compare CSV → {out_all}\n"); self.logbox.see(tk.END)
        except Exception:
            pass


def main() -> None:
    root = tk.Tk(); SimpleMASWGUI(root); root.mainloop()


if __name__ == "__main__":
    main()


