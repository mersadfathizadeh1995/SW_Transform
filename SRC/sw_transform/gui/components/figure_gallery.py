"""Figure gallery component for browsing and managing output images.

Copied from simple_app.py:
- lines 290-340 (Figures tab UI - toolbar, preview canvas, sidebar)
- lines 1224-1256 (refresh_gallery, _on_figure_selected, _load_preview_image)
- lines 1258-1321 (_render_fig_preview, zoom methods)
- lines 1380-1433 (_open_selected_figure, _delete_selected_figure, _build_ppt_from_gallery)
"""
from __future__ import annotations

import os
import tkinter as tk
from tkinter import ttk, messagebox, filedialog
from typing import Callable


class FigureGallery(tk.Frame):
    """Figure gallery with preview, zoom, and PPT export.
    
    Provides a split view with image preview on the left and
    file browser on the right. Supports zoom, pan, and PowerPoint
    export of all images.
    
    Usage:
        gallery = FigureGallery(parent, output_dir_var=output_folder_var)
        gallery.pack(fill="both", expand=True)
        
        gallery.refresh()
        gallery.open_selected()
        gallery.build_ppt()
    """
    
    def __init__(self, parent: tk.Widget,
                 output_dir_var: tk.StringVar | None = None,
                 icon_loader: Callable[[str, int], tk.PhotoImage | None] | None = None,
                 **kwargs):
        """Initialize the figure gallery.
        
        Args:
            parent: Parent widget
            output_dir_var: StringVar with output directory path
            icon_loader: Optional function to load icons
            **kwargs: Additional Frame options
        """
        super().__init__(parent, **kwargs)
        
        self.output_dir_var = output_dir_var or tk.StringVar()
        self.icon_loader = icon_loader
        
        # Preview state
        self._fig_image_pil = None
        self._fig_image_tk = None
        self._fig_scale = 1.0
        
        self._build_ui()
    
    def _build_ui(self):
        """Build the gallery UI.
        
        Copied from simple_app.py lines 290-340.
        """
        outer = tk.Frame(self)
        outer.pack(fill="both", expand=True)
        
        # Center preview area
        center = tk.Frame(outer)
        center.pack(side="left", fill="both", expand=True, padx=(6, 0), pady=6)
        
        # Toolbar
        toolbar = tk.Frame(center)
        toolbar.pack(fill="x")
        
        tk.Button(toolbar, text="−", width=3, 
                  command=lambda: self._fig_zoom_step(1/1.15)).pack(side="left")
        tk.Button(toolbar, text="+", width=3, 
                  command=lambda: self._fig_zoom_step(1.15)).pack(side="left", padx=(4, 0))
        tk.Button(toolbar, text="100%", 
                  command=self._fig_zoom_reset).pack(side="left", padx=(8, 0))
        
        tk.Label(toolbar, text="Fit:").pack(side="left", padx=(12, 2))
        self.fig_fit_mode = tk.StringVar(value="Auto")
        ttk.Combobox(
            toolbar,
            values=("Auto", "Width", "Height", "None"),
            width=8,
            state="readonly",
            textvariable=self.fig_fit_mode
        ).pack(side="left")
        
        self.fig_zoom_label = tk.StringVar(value="100%")
        tk.Label(toolbar, textvariable=self.fig_zoom_label).pack(side="left", padx=(10, 0))
        
        # Preview canvas
        prev_frame = tk.Frame(center)
        prev_frame.pack(fill="both", expand=True, pady=(4, 0))
        
        self.fig_prev_canvas = tk.Canvas(prev_frame, bg="#f3f3f3")
        vs = tk.Scrollbar(prev_frame, orient="vertical", command=self.fig_prev_canvas.yview)
        hs = tk.Scrollbar(prev_frame, orient="horizontal", command=self.fig_prev_canvas.xview)
        self.fig_prev_canvas.configure(yscrollcommand=vs.set, xscrollcommand=hs.set)
        
        self.fig_prev_canvas.grid(row=0, column=0, sticky="nsew")
        vs.grid(row=0, column=1, sticky="ns")
        hs.grid(row=1, column=0, sticky="ew")
        prev_frame.rowconfigure(0, weight=1)
        prev_frame.columnconfigure(0, weight=1)
        
        # Pan and zoom bindings
        self.fig_prev_canvas.bind("<ButtonPress-1>", 
            lambda e: self.fig_prev_canvas.scan_mark(e.x, e.y))
        self.fig_prev_canvas.bind("<B1-Motion>", 
            lambda e: self.fig_prev_canvas.scan_dragto(e.x, e.y, gain=1))
        self.fig_prev_canvas.bind("<MouseWheel>", self._on_fig_wheel_zoom)
        self.fig_prev_canvas.bind("<Configure>", lambda e: self._render_fig_preview())
        
        # Right sidebar explorer
        right = tk.Frame(outer, width=320)
        right.pack(side="left", fill="y", padx=6, pady=6)
        
        # Toolbar buttons
        sb_tools = tk.Frame(right)
        sb_tools.pack(fill="x")
        
        tk.Button(sb_tools, text="Refresh", command=self.refresh, pady=2).pack(side="left")
        tk.Button(sb_tools, text="Open", command=self.open_selected, pady=2).pack(side="left", padx=(6, 0))
        tk.Button(sb_tools, text="Delete", command=self.delete_selected, pady=2).pack(side="left", padx=(6, 0))
        
        btn_ppt = tk.Button(sb_tools, text=" PPT", command=self.build_ppt, 
                            compound="left", padx=6, pady=2)
        if self.icon_loader:
            ico_ppt = self.icon_loader("ic_ppt.png", 28)
            if ico_ppt is not None:
                btn_ppt.config(image=ico_ppt)
                btn_ppt._icon = ico_ppt
        btn_ppt.pack(side="left", padx=(6, 0))
        
        # File list treeview
        self.fig_list = ttk.Treeview(right, columns=("name", "path"), show="headings", height=22)
        self.fig_list.heading("name", text="Name")
        self.fig_list.column("name", width=180)
        self.fig_list.heading("path", text="Path")
        self.fig_list.column("path", width=120)
        self.fig_list.pack(fill="both", expand=True, pady=(6, 0))
        self.fig_list.bind("<<TreeviewSelect>>", self._on_figure_selected)
    
    def refresh(self):
        """Refresh the gallery from output folder.
        
        Copied from simple_app.py lines 1224-1238.
        """
        output_folder = self.output_dir_var.get()
        if not output_folder:
            messagebox.showwarning("Figures", "Select an output folder first.")
            return
        
        try:
            import glob
            self.fig_list.delete(*self.fig_list.get_children())
            for path in glob.glob(os.path.join(output_folder, "**", "*.png"), recursive=True):
                name = os.path.basename(path)
                self.fig_list.insert("", "end", values=(name, path))
        except Exception as e:
            messagebox.showerror("Figures", str(e))
    
    def _on_figure_selected(self, _):
        """Handle figure selection in treeview.
        
        Copied from simple_app.py lines 1240-1247.
        """
        sel = self.fig_list.selection()
        if not sel:
            return
        path = self.fig_list.item(sel[0], "values")[1]
        self._load_preview_image(path)
        self._render_fig_preview()
    
    def _load_preview_image(self, path: str):
        """Load an image for preview.
        
        Copied from simple_app.py lines 1249-1256.
        """
        try:
            from PIL import Image, ImageTk
        except Exception:
            self._fig_image_pil = None
            self._fig_image_tk = None
            return
        
        try:
            self._fig_image_pil = Image.open(path)
            self._fig_scale = 1.0
            self.fig_fit_mode.set("Auto")
        except Exception:
            self._fig_image_pil = None
            self._fig_image_tk = None
    
    def _render_fig_preview(self):
        """Render the preview image with current zoom/fit settings.
        
        Copied from simple_app.py lines 1258-1289.
        """
        if self._fig_image_pil is None:
            self.fig_prev_canvas.delete("all")
            self.fig_prev_canvas.config(scrollregion=(0, 0, 0, 0))
            return
        
        try:
            from PIL import ImageTk
            
            cw = self.fig_prev_canvas.winfo_width()
            ch = self.fig_prev_canvas.winfo_height()
            if cw <= 1 or ch <= 1:
                self.winfo_toplevel().update_idletasks()
                cw = self.fig_prev_canvas.winfo_width() or 900
                ch = self.fig_prev_canvas.winfo_height() or 600
            
            ow, oh = self._fig_image_pil.size
            mode = (self.fig_fit_mode.get() or "Auto").strip()
            
            if mode == "Width":
                scale = max(1e-3, (cw - 8) / ow)
            elif mode == "Height":
                scale = max(1e-3, (ch - 8) / oh)
            elif mode == "Auto":
                scale = max(1e-3, min((cw - 8) / ow, (ch - 8) / oh))
            else:
                scale = max(self._fig_scale, 1e-3)
            
            tw = max(50, int(ow * scale))
            th = max(50, int(oh * scale))
            
            self._fig_image_tk = ImageTk.PhotoImage(
                self._fig_image_pil.resize((tw, th))
            )
            
            self.fig_prev_canvas.delete("all")
            self.fig_prev_canvas.create_image(0, 0, image=self._fig_image_tk, anchor="nw")
            self.fig_prev_canvas.config(scrollregion=(0, 0, tw, th))
            
            try:
                self.fig_zoom_label.set(f"{int(round(scale * 100))}%")
            except Exception:
                pass
        except Exception:
            pass
    
    def _fig_zoom_step(self, factor: float):
        """Zoom by a factor.
        
        Copied from simple_app.py lines 1297-1303.
        """
        self.fig_fit_mode.set("None")
        self._fig_scale *= float(factor)
        if self._fig_scale < 0.05:
            self._fig_scale = 0.05
        self._render_fig_preview()
    
    def _fig_zoom_reset(self):
        """Reset zoom to auto-fit.
        
        Copied from simple_app.py lines 1305-1309.
        """
        self._fig_scale = 1.0
        self.fig_fit_mode.set("Auto")
        self._render_fig_preview()
    
    def _on_fig_wheel_zoom(self, event):
        """Handle mouse wheel zoom.
        
        Copied from simple_app.py lines 1311-1321.
        """
        step = 1.05
        self.fig_fit_mode.set("None")
        if event.delta > 0:
            self._fig_scale *= step
        else:
            self._fig_scale /= step
            if self._fig_scale < 0.05:
                self._fig_scale = 0.05
        self._render_fig_preview()
    
    def open_selected(self):
        """Open selected figure in system viewer.
        
        Copied from simple_app.py lines 1380-1393.
        """
        sel = self.fig_list.selection()
        if not sel:
            return
        path = self.fig_list.item(sel[0], "values")[1]
        try:
            import subprocess, sys
            if os.name == 'nt':
                os.startfile(path)
            elif sys.platform == 'darwin':
                subprocess.call(['open', path])
            else:
                subprocess.call(['xdg-open', path])
        except Exception:
            pass
    
    def delete_selected(self):
        """Delete selected figure.
        
        Copied from simple_app.py lines 1395-1403.
        """
        sel = self.fig_list.selection()
        if not sel:
            return
        path = self.fig_list.item(sel[0], "values")[1]
        try:
            os.remove(path)
            self.refresh()
        except Exception:
            pass
    
    def build_ppt(self):
        """Build PowerPoint from all images in gallery.
        
        Copied from simple_app.py lines 1405-1433.
        """
        output_folder = self.output_dir_var.get()
        if not output_folder:
            messagebox.showwarning("PowerPoint", "Select an output folder first.")
            return
        
        try:
            from pptx import Presentation
            from pptx.util import Inches
        except Exception:
            messagebox.showwarning("PowerPoint", 
                "python-pptx not installed. Run: pip install python-pptx")
            return
        
        try:
            import glob
            pngs = glob.glob(os.path.join(output_folder, "**", "*.png"), recursive=True)
            if not pngs:
                messagebox.showinfo("PowerPoint", "No PNG files found in output folder.")
                return
            
            out_ppt = os.path.join(output_folder, "slides_all_outputs.pptx")
            prs = Presentation()
            blank = prs.slide_layouts[6]
            
            sw_in = float(prs.slide_width) / 914400.0
            sh_in = float(prs.slide_height) / 914400.0
            tw_in = max(0.5, sw_in - 1.0)
            th_in = max(0.5, sh_in - 1.0)
            
            for img in pngs:
                slide = prs.slides.add_slide(blank)
                try:
                    from PIL import Image
                    Image.MAX_IMAGE_PIXELS = None
                    with Image.open(img) as im:
                        iw, ih = im.size
                    ar = iw / ih
                    target_ar = tw_in / max(th_in, 1e-6)
                    if target_ar > ar:
                        h_in = th_in
                        w_in = h_in * ar
                    else:
                        w_in = tw_in
                        h_in = w_in / max(ar, 1e-9)
                    x = Inches(max(0.0, (sw_in - w_in) / 2.0))
                    y = Inches(max(0.0, (sh_in - h_in) / 2.0))
                    slide.shapes.add_picture(img, x, y, width=Inches(w_in), height=Inches(h_in))
                except Exception:
                    pass
            
            prs.save(out_ppt)
            messagebox.showinfo("PowerPoint", f"Created: {out_ppt}")
            self.refresh()
        except Exception as e:
            messagebox.showerror("PowerPoint", str(e))
