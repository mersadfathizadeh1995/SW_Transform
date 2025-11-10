from __future__ import annotations

import os
import tkinter as tk
from tkinter import ttk, filedialog, messagebox

from sw_transform.core.service import run_single as svc_run_single, run_compare as svc_run_compare
from sw_transform.processing.registry import METHODS, dyn as dyn_func, compute_reverse_flag
from sw_transform.io.file_assignment import assign_files as assign_files_from_names


class SimpleMASWGUI:
    def __init__(self, root: tk.Tk):
        self.root = root
        self.root.title("MASW – SW_Transform GUI")
        try:
            self.root.geometry("1000x680"); self.root.minsize(900, 560)
        except Exception:
            pass
        # App icon (white background)
        try:
            # Titlebar/taskbar icon: prefer .ico; build from the largest PNG if needed
            base_assets = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "..", "..", "assets_big"))
            # search for the largest icon_app*.png
            src_png = None; max_size = -1
            if os.path.isdir(base_assets):
                for fn in os.listdir(base_assets):
                    if fn.lower().startswith("icon_app") and fn.lower().endswith(".png"):
                        p = os.path.join(base_assets, fn)
                        try:
                            from PIL import Image
                            with Image.open(p) as im:
                                w, h = im.size
                                if w*h > max_size:
                                    max_size = w*h; src_png = p
                        except Exception:
                            continue
            if src_png is not None:
                try:
                    from PIL import Image
                    import tempfile
                    ico_tmp = os.path.join(tempfile.gettempdir(), "sw_transform_icon_tmp.ico")
                    with Image.open(src_png) as im:
                        # build ICO with multiple sizes for better scaling
                        sizes = [(16,16),(24,24),(32,32),(48,48),(64,64),(128,128),(256,256)]
                        im.save(ico_tmp, format='ICO', sizes=sizes)
                    try:
                        self.root.iconbitmap(ico_tmp)
                    except Exception:
                        # fallback to iconphoto PNG
                        self._app_icon = tk.PhotoImage(file=src_png)
                        self.root.iconphoto(True, self._app_icon)
                except Exception:
                    pass
        except Exception:
            pass

        # icon cache
        self._icons: dict[str, tk.PhotoImage] = {}

        # state
        self.file_list: list[str] = []
        self.offsets: dict[str, str] = {}
        self.reverse_flags: dict[str, bool] = {}
        self.output_folder: str = ""
        self.method_key = tk.StringVar(value="fk")
        self.vibrosis_mode = tk.BooleanVar(value=False)  # False=hammer, True=vibrosis

        # inputs
        self.vmin_var = tk.StringVar(value="0"); self.vmax_var = tk.StringVar(value="5000")
        self.fmin_var = tk.StringVar(value="0"); self.fmax_var = tk.StringVar(value="100")
        self.time_start_var = tk.StringVar(value="0.0"); self.time_end_var = tk.StringVar(value="1.0")
        self.downsample_var = tk.BooleanVar(value=True); self.down_factor_var = tk.StringVar(value="16")
        self.numf_var = tk.StringVar(value="4000")
        self.grid_fk_var = tk.StringVar(value="4000"); self.tol_fk_var = tk.StringVar(value="0")
        self.grid_ps_var = tk.StringVar(value="1200"); self.vspace_ps_var = tk.StringVar(value="log"); self.tol_ps_var = tk.StringVar(value="0")
        self.dpi_var = tk.StringVar(value="200")
        self.figure_topic_var = tk.StringVar(value="3-D Dispersion (Freq vs. Velocity)")
        self.ppt_var = tk.BooleanVar(value=False)

        self._build_menu()
        self._build_ui()

    # UI
    def _build_menu(self):
        m = tk.Menu(self.root)
        filem = tk.Menu(m, tearoff=0)
        filem.add_command(label="Open SEG-2...", command=self.select_files)
        filem.add_command(label="Select Output Folder...", command=self.select_out)
        filem.add_separator(); filem.add_command(label="Exit", command=self.root.quit)
        m.add_cascade(label="File", menu=filem)

        datam = tk.Menu(m, tearoff=0)
        datam.add_command(label="Auto Assign Offsets (from filenames)", command=self._auto_assign_from_filenames)
        m.add_cascade(label="Data", menu=datam)

        self.root.config(menu=m)

    def _build_ui(self):
        left = tk.Frame(self.root, width=320); left.pack(side="left", fill="y")
        # Left toolbar with icon
        btn_open = tk.Button(left, text=" Open SEG-2...", command=self.select_files, compound="left", padx=6)
        ico = self._load_icon("ic_open.png", 56)
        if ico is not None:
            btn_open.config(image=ico)
        btn_open.pack(anchor="w", padx=8, pady=6)
        self.tree = ttk.Treeview(left, columns=("file","offset","rev"), show="headings", height=24)
        for col, w in zip(("file","offset","rev"), (220, 70, 40)):
            self.tree.heading(col, text=col.capitalize()); self.tree.column(col, width=w)
        self.tree.pack(fill="y", padx=8, pady=4, expand=True)
        self.tree.bind("<Double-1>", self._edit_cell)

        center = tk.Frame(self.root); center.pack(side="left", fill="both", expand=True)
        nb = ttk.Notebook(center)
        self.tab_inputs = tk.Frame(nb); self.tab_run = tk.Frame(nb); self.tab_fig = tk.Frame(nb)
        nb.add(self.tab_inputs, text="Inputs"); nb.add(self.tab_run, text="Run"); nb.add(self.tab_fig, text="Figures")
        nb.pack(fill="both", expand=True)

        # Inputs
        p = self.tab_inputs
        row = tk.Frame(p); row.pack(fill="x", padx=6, pady=4)
        tk.Label(row, text="Output folder:").pack(side="left")
        self.out_var = tk.StringVar(value="(not set)")
        tk.Label(row, textvariable=self.out_var, anchor="w").pack(side="left", fill="x", expand=True, padx=6)
        tk.Button(row, text="Select", command=self.select_out).pack(side="left")

        box = tk.LabelFrame(p, text="Picker limits"); box.pack(fill="x", padx=6, pady=4)
        for lab, var in (("Vmin", self.vmin_var), ("Vmax", self.vmax_var), ("Fmin", self.fmin_var), ("Fmax", self.fmax_var)):
            tk.Label(box, text=lab).pack(side="left"); tk.Entry(box, width=6, textvariable=var).pack(side="left", padx=3)

        tw = tk.LabelFrame(p, text="Time window & sampling"); tw.pack(fill="x", padx=6, pady=4)
        tk.Label(tw, text="Start (s)").pack(side="left"); tk.Entry(tw, width=6, textvariable=self.time_start_var).pack(side="left", padx=3)
        tk.Label(tw, text="End (s)").pack(side="left"); tk.Entry(tw, width=6, textvariable=self.time_end_var).pack(side="left", padx=3)
        tk.Checkbutton(tw, text="Downsample", variable=self.downsample_var).pack(side="left", padx=6)
        tk.Label(tw, text="Factor").pack(side="left"); tk.Entry(tw, width=4, textvariable=self.down_factor_var).pack(side="left", padx=3)
        tk.Label(tw, text="numf").pack(side="left"); tk.Entry(tw, width=6, textvariable=self.numf_var).pack(side="left", padx=3)

        pm = tk.LabelFrame(p, text="Per-method settings"); pm.pack(fill="x", padx=6, pady=4)
        fk = tk.Frame(pm); fk.pack(fill="x", pady=2)
        tk.Label(fk, text="FK/FDBF N").pack(side="left"); tk.Entry(fk, width=6, textvariable=self.grid_fk_var).pack(side="left", padx=4)
        tk.Label(fk, text="tol").pack(side="left"); tk.Entry(fk, width=6, textvariable=self.tol_fk_var).pack(side="left", padx=4)
        # Vibrosis checkbox on the right side
        tk.Checkbutton(fk, text="☐ Vibrosis (FDBF)", variable=self.vibrosis_mode).pack(side="right", padx=10)
        ps = tk.Frame(pm); ps.pack(fill="x", pady=2)
        tk.Label(ps, text="PS/SS N").pack(side="left"); tk.Entry(ps, width=6, textvariable=self.grid_ps_var).pack(side="left", padx=4)
        tk.Label(ps, text="vspace").pack(side="left")
        ttk.Combobox(ps, values=("linear","log"), width=6, textvariable=self.vspace_ps_var, state="readonly").pack(side="left", padx=4)
        tk.Label(ps, text="tol").pack(side="left"); tk.Entry(ps, width=6, textvariable=self.tol_ps_var).pack(side="left", padx=4)

        figbox = tk.LabelFrame(p, text="Figure / Export"); figbox.pack(fill="x", padx=6, pady=4)
        tk.Label(figbox, text="Figure DPI").pack(side="left"); tk.Entry(figbox, width=6, textvariable=self.dpi_var).pack(side="left", padx=4)
        tk.Label(figbox, text="Topic").pack(side="left", padx=(12,2)); tk.Entry(figbox, width=36, textvariable=self.figure_topic_var).pack(side="left", padx=2)

        # Array preview (embedded)
        arr_box = tk.LabelFrame(p, text="Array preview (embedded)")
        arr_box.pack(fill="both", expand=True, padx=6, pady=6)
        topbar = tk.Frame(arr_box); topbar.pack(fill="x", pady=(2,4))
        tk.Button(topbar, text="Preview Array / Waterfall", command=self.preview_array).pack(side="left")
        tk.Label(topbar, text="Display time (s):").pack(side="left", padx=(10,2))
        self.display_time_var = tk.StringVar(value="")
        tk.Entry(topbar, width=6, textvariable=self.display_time_var).pack(side="left")
        self.prev_host = tk.Frame(arr_box, height=300, bg="#f7f7f7"); self.prev_host.pack(fill="both", expand=True)
        self.prev_canvas_widget = None

        # Run tab
        r = self.tab_run
        top = tk.Frame(r); top.pack(pady=6)
        tk.Label(top, text="Transform:").pack(side="left")
        cmb = ttk.Combobox(top, values=[f"{k} – {d['label']}" for k,d in METHODS.items()], state="readonly", width=48)
        cmb.current(0); cmb.pack(side="left", padx=6)
        cmb.bind("<<ComboboxSelected>>", lambda ev: self.method_key.set(ev.widget.get().split(" –")[0]))
        row = tk.Frame(r); row.pack(pady=4)
        btn_run_sel = tk.Button(row, text=" Run Selected", command=lambda: self.run_single_processing(selected_only=True), compound="left", padx=6)
        btn_run_all = tk.Button(row, text=" Run All", command=lambda: self.run_single_processing(selected_only=False), compound="left", padx=6)
        ico_run = self._load_icon("ic_run.png", 56)
        if ico_run is not None:
            btn_run_sel.config(image=ico_run); btn_run_all.config(image=ico_run)
        btn_run_sel.pack(side="left", padx=4); btn_run_all.pack(side="left", padx=4)
        row2 = tk.Frame(r); row2.pack(pady=4)
        btn_cmp_sel = tk.Button(row2, text=" Compare Selected", command=lambda: self.run_compare_processing(selected_only=True), compound="left", padx=6)
        btn_cmp_all = tk.Button(row2, text=" Compare All", command=lambda: self.run_compare_processing(selected_only=False), compound="left", padx=6)
        ico_cmp = self._load_icon("ic_compare.png", 56)
        if ico_cmp is not None:
            btn_cmp_sel.config(image=ico_cmp); btn_cmp_all.config(image=ico_cmp)
        btn_cmp_sel.pack(side="left", padx=4); btn_cmp_all.pack(side="left", padx=4)
        opt = tk.Frame(r); opt.pack(pady=2)
        tk.Checkbutton(opt, text="Create PowerPoint after run", variable=self.ppt_var).pack(side="left")
        # Log box
        try:
            from tkinter import scrolledtext
            self.logbox = scrolledtext.ScrolledText(r, width=92, height=10)
            self.logbox.pack(fill="both", expand=True, padx=6, pady=6)
        except Exception:
            self.logbox = None

        # Bottom-right progress bar
        pb_row = tk.Frame(r); pb_row.pack(fill="x", padx=6, pady=(0,6))
        self.pb = ttk.Progressbar(pb_row, orient="horizontal", mode="determinate", length=220)
        self.pb.pack(side="right")
        self.pb_label = tk.Label(pb_row, text="Idle", anchor="e")
        self.pb_label.pack(side="right", padx=(0,8))

        # Figures tab UI (split: center preview + right sidebar explorer)
        ftab = self.tab_fig
        outer = tk.Frame(ftab); outer.pack(fill="both", expand=True)
        # Center preview area
        center = tk.Frame(outer); center.pack(side="left", fill="both", expand=True, padx=(6,0), pady=6)
        toolbar = tk.Frame(center); toolbar.pack(fill="x")
        tk.Button(toolbar, text="−", width=3, command=lambda: self._fig_zoom_step(1/1.15)).pack(side="left")
        tk.Button(toolbar, text="+", width=3, command=lambda: self._fig_zoom_step(1.15)).pack(side="left", padx=(4,0))
        tk.Button(toolbar, text="100%", command=self._fig_zoom_reset).pack(side="left", padx=(8,0))
        tk.Label(toolbar, text="Fit:").pack(side="left", padx=(12,2))
        self.fig_fit_mode = tk.StringVar(value="Auto")
        ttk.Combobox(toolbar, values=("Auto","Width","Height","None"), width=8, state="readonly", textvariable=self.fig_fit_mode).pack(side="left")
        self.fig_zoom_label = tk.StringVar(value="100%")
        tk.Label(toolbar, textvariable=self.fig_zoom_label).pack(side="left", padx=(10,0))
        # Preview canvas
        prev_frame = tk.Frame(center); prev_frame.pack(fill="both", expand=True, pady=(4,0))
        self.fig_prev_canvas = tk.Canvas(prev_frame, bg="#f3f3f3")
        vs = tk.Scrollbar(prev_frame, orient="vertical", command=self.fig_prev_canvas.yview)
        hs = tk.Scrollbar(prev_frame, orient="horizontal", command=self.fig_prev_canvas.xview)
        self.fig_prev_canvas.configure(yscrollcommand=vs.set, xscrollcommand=hs.set)
        self.fig_prev_canvas.grid(row=0, column=0, sticky="nsew"); vs.grid(row=0, column=1, sticky="ns"); hs.grid(row=1, column=0, sticky="ew")
        prev_frame.rowconfigure(0, weight=1); prev_frame.columnconfigure(0, weight=1)
        self.fig_prev_canvas.bind("<ButtonPress-1>", lambda e: self.fig_prev_canvas.scan_mark(e.x, e.y))
        self.fig_prev_canvas.bind("<B1-Motion>", lambda e: self.fig_prev_canvas.scan_dragto(e.x, e.y, gain=1))
        self.fig_prev_canvas.bind("<MouseWheel>", self._on_fig_wheel_zoom)
        self.fig_prev_canvas.bind("<Configure>", lambda e: self._render_fig_preview())

        # Right sidebar explorer
        right = tk.Frame(outer, width=320); right.pack(side="left", fill="y", padx=6, pady=6)
        sb_tools = tk.Frame(right); sb_tools.pack(fill="x")
        tk.Button(sb_tools, text="Refresh", command=self.refresh_gallery).pack(side="left")
        tk.Button(sb_tools, text="Open", command=self._open_selected_figure).pack(side="left", padx=(6,0))
        tk.Button(sb_tools, text="Delete", command=self._delete_selected_figure).pack(side="left", padx=(6,0))
        btn_ppt = tk.Button(sb_tools, text=" Create PPT", command=self._build_ppt_from_gallery, compound="left", padx=6)
        ico_ppt = self._load_icon("ic_ppt.png", 56)
        if ico_ppt is not None:
            btn_ppt.config(image=ico_ppt)
        btn_ppt.pack(side="left", padx=(6,0))
        self.fig_list = ttk.Treeview(right, columns=("name","path"), show="headings", height=22)
        self.fig_list.heading("name", text="Name"); self.fig_list.column("name", width=180)
        self.fig_list.heading("path", text="Path"); self.fig_list.column("path", width=120)
        self.fig_list.pack(fill="both", expand=True, pady=(6,0))
        self.fig_list.bind("<<TreeviewSelect>>", self._on_figure_selected)

        # Preview state
        self._fig_image_pil = None
        self._fig_image_tk = None
        self._fig_scale = 1.0

    # actions
    def select_files(self):
        files = filedialog.askopenfilenames(title="Select SEG-2 Files", filetypes=[("SEG-2 .dat","*.dat"),("All files","*.*")])
        if not files:
            return
        self.file_list = list(files)
        bases = [os.path.splitext(os.path.basename(f))[0] for f in files]
        self.offsets = {b: "+0" for b in bases}; self.reverse_flags = {b: False for b in bases}
        self.tree.delete(*self.tree.get_children())
        for b in bases:
            self.tree.insert("", "end", values=(b, "+0", "☐"))

    def select_out(self):
        folder = filedialog.askdirectory(title="Choose output folder")
        if folder:
            self.output_folder = folder; self.out_var.set(folder)

    def _edit_cell(self, ev):
        item = self.tree.identify_row(ev.y); col = self.tree.identify_column(ev.x)
        if not item: return
        col_idx = int(col.lstrip("#")); vals = list(self.tree.item(item, "values")); base = vals[0]
        x, y, w, h = self.tree.bbox(item, col)
        if col_idx == 2:
            e = tk.Entry(self.root); e.insert(0, vals[1])
            e.place(x=x, y=y + self.tree.winfo_rooty()-self.root.winfo_rooty(), width=w, height=h); e.focus()
            def _done(_):
                vals[1] = e.get(); self.offsets[base] = vals[1]; self.tree.item(item, values=vals); e.destroy()
            e.bind("<Return>", _done); e.bind("<FocusOut>", _done)
        elif col_idx == 3:
            flag = vals[2] != "☑"; self.reverse_flags[base] = flag; vals[2] = "☑" if flag else "☐"; self.tree.item(item, values=vals)

    def _collect_common(self):
        try:
            pick_vmin = float(self.vmin_var.get()); pick_vmax = float(self.vmax_var.get())
            pick_fmin = float(self.fmin_var.get()); pick_fmax = float(self.fmax_var.get())
        except ValueError:
            raise ValueError("Invalid pick/plot limits")
        try:
            st = float(self.time_start_var.get()); en = float(self.time_end_var.get())
            downsample = bool(self.downsample_var.get()); dfac = int(self.down_factor_var.get()); numf = int(self.numf_var.get())
        except Exception as e:
            raise ValueError(f"Invalid sampling: {e}")
        try:
            fig_dpi = int(self.dpi_var.get()); assert fig_dpi > 0
        except Exception:
            raise ValueError("Figure DPI must be a positive integer")
        return pick_vmin, pick_vmax, pick_fmin, pick_fmax, st, en, downsample, dfac, numf, fig_dpi

    def _resolve_selected_path(self) -> str | None:
        if not self.file_list:
            return None
        sel = self.tree.selection()
        if sel:
            want = self.tree.item(sel[0], "values")[0]
            for f in self.file_list:
                if os.path.splitext(os.path.basename(f))[0] == want:
                    return f
        return self.file_list[0]

    def preview_array(self):
        # Build an embedded preview (array schematic + waterfall)
        path = self._resolve_selected_path()
        if path is None:
            messagebox.showerror("Preview", "No file selected."); return
        try:
            from matplotlib.figure import Figure
            from matplotlib.backends.backend_tkagg import FigureCanvasTkAgg
            from sw_transform.processing.seg2 import load_seg2_ar
            from sw_transform.processing.preprocess import preprocess_data
            import numpy as np
            time, T, Shotpoint, Spacing, dt, _ = load_seg2_ar(path)
            st = float(self.time_start_var.get()); en = float(self.time_end_var.get())
            ds = bool(self.downsample_var.get()); df = int(self.down_factor_var.get()); nf = int(self.numf_var.get())
            Tpre, time_pre, dt2 = preprocess_data(T, time, dt, reverse_shot=False, start_time=st, end_time=en, do_downsample=ds, down_factor=df, numf=nf)
            # Optional display time trim
            disp_txt = (self.display_time_var.get() or "").strip()
            if disp_txt:
                try:
                    disp_time = float(disp_txt)
                    if disp_time > 0:
                        n_keep = int(np.clip(disp_time / dt2, 1, Tpre.shape[0]))
                        Tpre = Tpre[:n_keep, :]; time_pre = time_pre[:n_keep]
                except Exception:
                    pass
            positions = np.arange(Tpre.shape[1], dtype=float) * float(Spacing)
            fig = Figure(figsize=(7.5, 6.0), dpi=100)
            gs = fig.add_gridspec(2, 1, height_ratios=[1, 3], hspace=0.42)
            ax1 = fig.add_subplot(gs[0]); ax2 = fig.add_subplot(gs[1])
            # schematic
            ax1.plot(positions, np.zeros_like(positions), "^", color="green", label="Sensor")
            # Source position: prefer table override
            try:
                base = os.path.splitext(os.path.basename(path))[0]
                off_txt = (self.offsets.get(base, "+0") or "+0").strip().replace("m", "")
                if off_txt.startswith("+"):
                    src_x = float(off_txt[1:])
                else:
                    src_x = float(off_txt)
            except Exception:
                src_x = float(Shotpoint)
            ax1.plot([src_x], [0.0], "D", color="tab:red", label="Source")
            ax1.set_yticks([]); ax1.set_xlabel("Distance (m)"); ax1.legend(loc="upper right"); ax1.set_title("Array schematic")
            # waterfall
            traces = Tpre.copy().T
            denom = np.max(np.abs(traces), axis=1, keepdims=True); denom[denom==0]=1.0; traces = traces/denom
            spacing = float(np.mean(np.diff(positions))) if len(positions)>1 else 1.0
            scale = 0.5 * spacing
            for tr, x0 in zip(traces, positions):
                ax2.plot(tr*scale + x0, time_pre, color="b", linewidth=0.5)
            ax2.invert_yaxis(); ax2.set_xlabel("Distance (m)"); ax2.set_ylabel("Time (s)"); ax2.set_title("Waterfall (normalized)")
            # mount
            if self.prev_canvas_widget is not None:
                try: self.prev_canvas_widget.destroy()
                except Exception: pass
            canvas = FigureCanvasTkAgg(fig, master=self.prev_host)
            self.prev_canvas_widget = canvas.get_tk_widget(); self.prev_canvas_widget.pack(fill="both", expand=True)
            canvas.draw()
        except Exception as e:
            messagebox.showerror("Preview", str(e))

    def run_single_processing(self, selected_only: bool=False):
        if not self.file_list:
            messagebox.showerror("No files", "Select SEG-2 files first."); return
        if not self.output_folder:
            messagebox.showerror("No output", "Choose an output folder."); return
        key = self.method_key.get()
        if key in ("fk","fdbf"):
            grid_n = int(self.grid_fk_var.get()); tol = float(self.tol_fk_var.get()); vspace = "linear"
        else:
            grid_n = int(self.grid_ps_var.get()); tol = float(self.tol_ps_var.get()); vspace = self.vspace_ps_var.get() or "linear"
        pick_vmin, pick_vmax, pick_fmin, pick_fmax, st, en, downsample, dfac, numf, fig_dpi = self._collect_common()
        try:
            import matplotlib as mpl; old_dpi = mpl.rcParams.get('savefig.dpi', 'figure'); mpl.rcParams['savefig.dpi'] = fig_dpi
        except Exception:
            old_dpi = None
        try:
            paths = list(self.file_list)
            if selected_only:
                sel = self.tree.selection(); selected_bases = {self.tree.item(i, "values")[0] for i in sel} if sel else set()
                paths = [p for p in self.file_list if os.path.splitext(os.path.basename(p))[0] in selected_bases]
            # prep progress
            total = max(1, len(paths)); self.pb.config(maximum=total, value=0); self.pb_label.config(text="Processing…")
            self.root.update_idletasks()
            success = 0; failures = 0
            for path in paths:
                base = os.path.splitext(os.path.basename(path))[0]
                offset = self.offsets.get(base, "+0")
                user_rev = self.reverse_flags.get(base, False)
                rev = compute_reverse_flag(bool(user_rev), key)
                # Convert vibrosis boolean to source_type string
                source_type = "vibrosis" if self.vibrosis_mode.get() else "hammer"
                params = dict(path=path, base=base, key=key, offset=offset, outdir=self.output_folder,
                              pick_vmin=pick_vmin, pick_vmax=pick_vmax, pick_fmin=pick_fmin, pick_fmax=pick_fmax,
                              st=st, en=en, downsample=downsample, dfac=dfac, numf=numf,
                              grid_n=grid_n, tol=tol, vspace=vspace, dpi=fig_dpi, rev=rev,
                              topic=(self.figure_topic_var.get() or ""),
                              source_type=source_type)
                if self.logbox: self.logbox.insert(tk.END, f"Running {base} [{key}]...\n"); self.logbox.see(tk.END)
                _b, ok, out = svc_run_single(params)
                if ok and isinstance(out, str) and out.lower().endswith('.png'):
                    success += 1
                    if self.logbox: self.logbox.insert(tk.END, f"Saved: {out}\n"); self.logbox.see(tk.END)
                else:
                    failures += 1
                    if self.logbox: self.logbox.insert(tk.END, f"Failed {base}: {out}\n"); self.logbox.see(tk.END)
                # progress update
                self.pb.step(1); self.root.update_idletasks()
            # After all files processed, create combined CSV
            if success > 0:
                try:
                    self._create_combined_csv_single_method(key, paths)
                    if self.logbox:
                        self.logbox.insert(tk.END, f"Created combined_{key}.csv\n")
                        self.logbox.see(tk.END)
                except Exception as e:
                    if self.logbox:
                        self.logbox.insert(tk.END, f"Warning: Could not create combined CSV: {e}\n")
                        self.logbox.see(tk.END)
            if success > 0 and failures == 0:
                messagebox.showinfo("Run", f"Completed ({success} file(s))")
            elif success > 0:
                messagebox.showwarning("Run", f"Completed with errors: {success} success, {failures} failed. See log.")
            else:
                messagebox.showerror("Run", "No outputs created. Check the log for errors.")
            # Optional PPT
            if success > 0 and bool(self.ppt_var.get()):
                self._build_ppt_from_gallery()
            # Build combined CSV across all shot offsets
            if success > 1:
                self._build_combined_csv_single(key, paths)
        except Exception as e:
            messagebox.showerror("Run", str(e))
        finally:
            try:
                import matplotlib as mpl
                if old_dpi is not None:
                    mpl.rcParams['savefig.dpi'] = old_dpi
            except Exception:
                pass
            self.pb_label.config(text="Idle"); self.pb.config(value=0)

    def run_compare_processing(self, selected_only: bool=False):
        if not self.file_list:
            messagebox.showerror("No files", "Select SEG-2 files first."); return
        if not self.output_folder:
            messagebox.showerror("No output", "Choose an output folder."); return
        pick_vmin, pick_vmax, pick_fmin, pick_fmax, st, en, downsample, dfac, numf, fig_dpi = self._collect_common()
        n_fk = int(self.grid_fk_var.get()); tol_fk = float(self.tol_fk_var.get())
        n_ps = int(self.grid_ps_var.get()); vspace_ps = self.vspace_ps_var.get() or "linear"
        paths = list(self.file_list)
        if selected_only:
            sel = self.tree.selection(); selected_bases = {self.tree.item(i, "values")[0] for i in sel} if sel else set()
            paths = [p for p in self.file_list if os.path.splitext(os.path.basename(p))[0] in selected_bases]
        success = 0; failures = 0
        total = max(1, len(paths)); self.pb.config(maximum=total, value=0); self.pb_label.config(text="Comparing…"); self.root.update_idletasks()
        for path in paths:
            base = os.path.splitext(os.path.basename(path))[0]
            offset = self.offsets.get(base, "+0")
            user_rev = self.reverse_flags.get(base, False)
            # Convert vibrosis boolean to source_type string
            source_type = "vibrosis" if self.vibrosis_mode.get() else "hammer"
            params = dict(path=path, base=base, outdir=self.output_folder, offset=offset,
                          pick_vmin=pick_vmin, pick_vmax=pick_vmax, pick_fmin=pick_fmin, pick_fmax=pick_fmax,
                          st=st, en=en, downsample=downsample, dfac=dfac, numf=numf,
                          n_fk=n_fk, tol_fk=tol_fk, n_ps=n_ps, vspace_ps=vspace_ps,
                          rev_fk=(not user_rev), rev_ps=(not user_rev), rev_fdbf=user_rev, rev_ss=user_rev,
                          topic=(self.figure_topic_var.get() or ""),
                          source_type=source_type)
            if self.logbox: self.logbox.insert(tk.END, f"Compare {base}...\n"); self.logbox.see(tk.END)
            _b, ok, out = svc_run_compare(params)
            if ok and isinstance(out, str) and out.lower().endswith('.png'):
                success += 1
                if self.logbox: self.logbox.insert(tk.END, f"Saved: {out}\n"); self.logbox.see(tk.END)
            else:
                failures += 1
                if self.logbox: self.logbox.insert(tk.END, f"Failed {base}: {out}\n"); self.logbox.see(tk.END)
            self.pb.step(1); self.root.update_idletasks()
        # After all files processed, create combined CSV for compare mode
        if success > 0:
            try:
                self._create_combined_csv_compare_mode(paths)
                if self.logbox:
                    self.logbox.insert(tk.END, "Created combined_compare_all.csv\n")
                    self.logbox.see(tk.END)
            except Exception as e:
                if self.logbox:
                    self.logbox.insert(tk.END, f"Warning: Could not create combined CSV: {e}\n")
                    self.logbox.see(tk.END)
        if success > 0 and failures == 0:
            messagebox.showinfo("Run", f"Comparison completed ({success} file(s))")
        elif success > 0:
            messagebox.showwarning("Run", f"Comparison completed with errors: {success} success, {failures} failed. See log.")
        else:
            messagebox.showerror("Run", "No comparison outputs created. Check the log for errors.")
        # Optional PPT
        if success > 0 and bool(self.ppt_var.get()):
            self._build_ppt_from_gallery()
        # Build global combined CSV across all files (4 methods)
        if success > 1:
            self._build_combined_csv_compare(paths)
        self.pb_label.config(text="Idle"); self.pb.config(value=0)

    def _create_combined_csv_single_method(self, key: str, paths: list):
        """Aggregate all per-shot CSVs into combined_<method>.csv."""
        import csv
        import glob
        # Find all per-shot CSVs for this method
        pattern = os.path.join(self.output_folder, f"*_{key}_*.csv")
        csv_files = glob.glob(pattern)
        # Exclude any compare or combined files
        csv_files = [f for f in csv_files if not ('_compare' in f or 'combined_' in os.path.basename(f))]
        if not csv_files:
            return
        # Read all CSVs and organize by (base, offset)
        data_map = {}  # (base, offset) -> (freq, vel, wav)
        for csv_path in csv_files:
            try:
                with open(csv_path, 'r', newline='') as f:
                    reader = csv.reader(f)
                    header = next(reader, None)
                    if not header or len(header) < 2:
                        continue
                    freq_list = []; vel_list = []; wav_list = []
                    for row in reader:
                        if len(row) >= 2:
                            freq_list.append(row[0])
                            vel_list.append(row[1])
                            wav_list.append(row[2] if len(row) > 2 else "")
                    # Extract base and offset from filename
                    fname = os.path.basename(csv_path)
                    # Format: <base>_<method>_<offset_tag>.csv
                    # e.g., "file1_fk_p66.csv" -> base="file1", offset="+66"
                    parts = fname.replace(".csv", "").split("_")
                    if len(parts) >= 3:
                        # Find the method position in the filename
                        method_idx = -1
                        for idx, part in enumerate(parts):
                            if part == key:
                                method_idx = idx
                                break
                        if method_idx >= 0 and method_idx < len(parts) - 1:
                            base = "_".join(parts[:method_idx])
                            offset_tag = parts[method_idx + 1]
                            # Convert offset_tag back to display format
                            if offset_tag.startswith('p'):
                                offset = "+" + offset_tag[1:] + "m"
                            elif offset_tag.startswith('m'):
                                offset = "-" + offset_tag[1:] + "m"
                            else:
                                offset = offset_tag
                            data_map[(base, offset)] = (freq_list, vel_list, wav_list)
            except Exception as e:
                if self.logbox:
                    self.logbox.insert(tk.END, f"Warning: Failed to read {csv_path}: {e}\n")
                continue
        if not data_map:
            return
        # Create combined CSV
        combined_path = os.path.join(self.output_folder, f"combined_{key}.csv")
        with open(combined_path, 'w', newline='') as f:
            writer = csv.writer(f)
            # Build header
            header = []
            sorted_keys = sorted(data_map.keys())
            for base, offset in sorted_keys:
                header.extend([f"freq({key}_{offset})", f"vel({key}_{offset})", f"wav({key}_{offset})"])
            writer.writerow(header)
            # Find max length
            max_len = max(len(data_map[k][0]) for k in sorted_keys) if sorted_keys else 0
            # Write rows
            for i in range(max_len):
                row = []
                for base, offset in sorted_keys:
                    freq_list, vel_list, wav_list = data_map[(base, offset)]
                    row.append(freq_list[i] if i < len(freq_list) else "")
                    row.append(vel_list[i] if i < len(vel_list) else "")
                    row.append(wav_list[i] if i < len(wav_list) else "")
                writer.writerow(row)

    def _create_combined_csv_compare_mode(self, paths: list):
        """Aggregate all per-file compare CSVs into combined_compare_all.csv."""
        import csv
        import glob
        # Find all per-file compare CSVs
        pattern = os.path.join(self.output_folder, "*_compare.csv")
        csv_files = glob.glob(pattern)
        if not csv_files:
            return
        # Read all CSVs
        data_map = {}  # base -> list of rows
        headers_map = {}  # base -> header
        for csv_path in csv_files:
            try:
                with open(csv_path, 'r', newline='') as f:
                    reader = csv.reader(f)
                    header = next(reader, None)
                    if not header:
                        continue
                    rows = list(reader)
                    base = os.path.basename(csv_path).replace("_compare.csv", "")
                    data_map[base] = rows
                    headers_map[base] = header
            except Exception:
                continue
        if not data_map:
            return
        # Create combined CSV
        combined_path = os.path.join(self.output_folder, "combined_compare_all.csv")
        with open(combined_path, 'w', newline='') as f:
            writer = csv.writer(f)
            # Build combined header
            sorted_bases = sorted(data_map.keys())
            combined_header = []
            for base in sorted_bases:
                combined_header.extend(headers_map.get(base, []))
            writer.writerow(combined_header)
            # Find max length
            max_len = max(len(data_map[b]) for b in sorted_bases) if sorted_bases else 0
            # Write rows
            for i in range(max_len):
                row = []
                for base in sorted_bases:
                    file_rows = data_map[base]
                    if i < len(file_rows):
                        row.extend(file_rows[i])
                    else:
                        # Pad with empty cells matching header width
                        row.extend([""] * len(headers_map.get(base, [])))
                writer.writerow(row)

    # data helpers
    def _auto_assign_from_filenames(self):
        if not self.file_list:
            messagebox.showerror("No files", "Select SEG-2 files first."); return
        try:
            rows = assign_files_from_names(self.file_list, recursive=False, include_unknown=True)
            new_offsets = {}; new_reverse = {}
            for r in rows:
                try:
                    fp = str(getattr(r, 'file_path', ''))
                    b = os.path.splitext(os.path.basename(fp))[0]
                    off = getattr(r, 'offset_m', None); rev = bool(getattr(r, 'reverse', False))
                    if off is not None:
                        val = float(off); sign = "+" if val >= 0 else ""; new_offsets[b] = f"{sign}{int(round(val))}"
                    new_reverse[b] = rev
                except Exception:
                    continue
            for item in self.tree.get_children():
                vals = list(self.tree.item(item, "values")); base = vals[0]
                if base in new_offsets:
                    vals[1] = new_offsets[base]; self.offsets[base] = vals[1]
                if base in new_reverse:
                    self.reverse_flags[base] = bool(new_reverse[base]); vals[2] = "☑" if self.reverse_flags[base] else "☐"
                self.tree.item(item, values=vals)
            messagebox.showinfo("Assign", "Offsets/reverse assigned from filenames.")
        except Exception as e:
            messagebox.showerror("Assign", str(e))

    # --- Figures tab helpers ---
    def refresh_gallery(self):
        if not self.output_folder:
            messagebox.showwarning("Figures", "Select an output folder first."); return
        try:
            import glob, os
            self.fig_list.delete(*self.fig_list.get_children())
            for path in glob.glob(os.path.join(self.output_folder, "**", "*.png"), recursive=True):
                name = os.path.basename(path)
                self.fig_list.insert("", "end", values=(name, path))
        except Exception as e:
            messagebox.showerror("Figures", str(e))

    def _on_figure_selected(self, _):
        sel = self.fig_list.selection()
        if not sel:
            return
        path = self.fig_list.item(sel[0], "values")[1]
        self._load_preview_image(path)
        self._render_fig_preview()

    def _load_preview_image(self, path: str):
        try:
            from PIL import Image, ImageTk
        except Exception:
            self._fig_image_pil = None; self._fig_image_tk = None; return
        try:
            self._fig_image_pil = Image.open(path)
            self._fig_scale = 1.0
            self.fig_fit_mode.set("Auto")
        except Exception:
            self._fig_image_pil = None; self._fig_image_tk = None

    def _render_fig_preview(self):
        if self._fig_image_pil is None:
            self.fig_prev_canvas.delete("all"); self.fig_prev_canvas.config(scrollregion=(0,0,0,0)); return
        try:
            from PIL import ImageTk
            cw = self.fig_prev_canvas.winfo_width(); ch = self.fig_prev_canvas.winfo_height()
            if cw <= 1 or ch <= 1:
                self.root.update_idletasks(); cw = self.fig_prev_canvas.winfo_width() or 900; ch = self.fig_prev_canvas.winfo_height() or 600
            ow, oh = self._fig_image_pil.size
            mode = (self.fig_fit_mode.get() or "Auto").strip()
            if mode == "Width":
                scale = max(1e-3, (cw-8)/ow)
            elif mode == "Height":
                scale = max(1e-3, (ch-8)/oh)
            elif mode == "Auto":
                scale = max(1e-3, min((cw-8)/ow, (ch-8)/oh))
            else:
                scale = max(self._fig_scale, 1e-3)
            tw = max(50, int(ow*scale)); th = max(50, int(oh*scale))
            self._fig_image_tk = ImageTk.PhotoImage(self._fig_image_pil.resize((tw, th)))
            self.fig_prev_canvas.delete("all")
            self.fig_prev_canvas.create_image(0, 0, image=self._fig_image_tk, anchor="nw")
            self.fig_prev_canvas.config(scrollregion=(0,0,tw,th))
            try:
                self.fig_zoom_label.set(f"{int(round(scale*100))}%")
            except Exception:
                pass
        except Exception:
            pass

    def _fig_zoom_step(self, factor: float):
        self.fig_fit_mode.set("None")
        self._fig_scale *= float(factor)
        if self._fig_scale < 0.05:
            self._fig_scale = 0.05
        self._render_fig_preview()

    def _fig_zoom_reset(self):
        self._fig_scale = 1.0
        self.fig_fit_mode.set("Auto")
        self._render_fig_preview()

    def _on_fig_wheel_zoom(self, event):
        step = 1.05
        self.fig_fit_mode.set("None")
        if event.delta > 0:
            self._fig_scale *= step
        else:
            self._fig_scale /= step
            if self._fig_scale < 0.05:
                self._fig_scale = 0.05
        self._render_fig_preview()

    # --- icon helper ---
    def _asset_path(self, name: str) -> str:
        """Locate asset file, handling both exact names and @NxN suffixed variants."""
        base = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "..", "..", "assets_big"))
        if not os.path.isdir(base):
            return os.path.join(base, name)  # return path anyway
        # Try exact match first
        p = os.path.join(base, name)
        if os.path.isfile(p):
            return p
        # fallback: pick the largest file that starts with the prefix
        prefix = os.path.splitext(name)[0]
        try:
            best = None; best_area = -1
            for fn in os.listdir(base):
                # Match prefix (e.g., "ic_open" matches "ic_open@180x180.png")
                if fn.lower().startswith(prefix.lower()) and fn.lower().endswith('.png'):
                    cand = os.path.join(base, fn)
                    try:
                        from PIL import Image
                        with Image.open(cand) as im:
                            w, h = im.size
                            if w*h > best_area:
                                best = cand; best_area = w*h
                    except Exception:
                        continue
            if best:
                return best
        except Exception:
            pass
        return p

    def _load_icon(self, name: str, size: int) -> tk.PhotoImage | None:
        """Load and cache an icon, auto-cropping transparent margins and centering."""
        try:
            key = f"{name}:{size}"
            if key in self._icons:
                return self._icons[key]
            from PIL import Image, ImageOps, ImageTk
            p = self._asset_path(name)
            if not os.path.isfile(p):
                return None
            im = Image.open(p).convert("RGBA")
            # Remove transparent margins
            # Auto-crop transparent margins and zoom to fill entire square (no padding, no margins)
            alpha = im.split()[-1]
            bbox = alpha.getbbox()
            if bbox:
                im = im.crop(bbox)
            from PIL import ImageOps
            im = ImageOps.fit(im, (size, size), method=Image.Resampling.LANCZOS)
            # Place on white background (full-bleed)
            bg = Image.new("RGBA", (size, size), (255, 255, 255, 255))
            bg.alpha_composite(im, dest=(0, 0))
            tkimg = ImageTk.PhotoImage(bg)
            self._icons[key] = tkimg
            return tkimg
        except Exception:
            return None

    def _open_selected_figure(self):
        sel = self.fig_list.selection()
        if not sel: return
        path = self.fig_list.item(sel[0], "values")[1]
        try:
            import subprocess, sys
            if os.name == 'nt':
                os.startfile(path)
            elif sys.platform == 'darwin':
                subprocess.call(['open', path])
            else:
                subprocess.call(['xdg-open', path])
        except Exception:
            pass

    def _delete_selected_figure(self):
        sel = self.fig_list.selection()
        if not sel: return
        path = self.fig_list.item(sel[0], "values")[1]
        try:
            os.remove(path)
            self.refresh_gallery()
        except Exception:
            pass

    def _build_ppt_from_gallery(self):
        if not self.output_folder:
            messagebox.showwarning("PowerPoint", "Select an output folder first."); return
        try:
            from pptx import Presentation
            from pptx.util import Inches
        except Exception:
            messagebox.showwarning("PowerPoint", "python-pptx not installed. Run: pip install python-pptx"); return
        try:
            import glob, os
            pngs = glob.glob(os.path.join(self.output_folder, "**", "*.png"), recursive=True)
            if not pngs:
                messagebox.showinfo("PowerPoint", "No PNG files found in output folder."); return
            out_ppt = os.path.join(self.output_folder, "slides_all_outputs.pptx")
            prs = Presentation()
            blank = prs.slide_layouts[6]
            sw_in = float(prs.slide_width) / 914400.0
            sh_in = float(prs.slide_height) / 914400.0
            tw_in = max(0.5, sw_in - 1.0)
            th_in = max(0.5, sh_in - 1.0)
            for img in pngs:
                slide = prs.slides.add_slide(blank)
                try:
                    from PIL import Image
                    Image.MAX_IMAGE_PIXELS = None
                    with Image.open(img) as im: iw, ih = im.size
                    ar = iw/ih; target_ar = tw_in/max(th_in, 1e-6)
                    if target_ar > ar:
                        h_in = th_in; w_in = h_in*ar
                    else:
                        w_in = tw_in; h_in = w_in/max(ar, 1e-9)
                    x = Inches(max(0.0, (sw_in - w_in) / 2.0)); y = Inches(max(0.0, (sh_in - h_in) / 2.0))
                    slide.shapes.add_picture(img, x, y, width=Inches(w_in), height=Inches(h_in))
                except Exception:
                    pass
            prs.save(out_ppt)
            messagebox.showinfo("PowerPoint", f"Created: {out_ppt}")
            self.refresh_gallery()
        except Exception as e:
            messagebox.showerror("PowerPoint", str(e))

    def _build_combined_csv_single(self, key: str, paths: list[str]):
        """Aggregate per-shot CSVs into combined_<method>.csv across all offsets."""
        try:
            import csv
            combined_rows = []
            for path in paths:
                base = os.path.splitext(os.path.basename(path))[0]
                offset = self.offsets.get(base, "+0")
                _off = str(offset).strip().replace(" ", "").replace("m", "")
                if _off.startswith("+"):
                    _off_tag = "p" + _off[1:]
                elif _off.startswith("-"):
                    _off_tag = "m" + _off[1:]
                else:
                    _off_tag = _off
                per_csv = os.path.join(self.output_folder, f"{base}_{key}_{_off_tag}.csv")
                if not os.path.isfile(per_csv):
                    continue
                frq, vel, wl = [], [], []
                try:
                    with open(per_csv, "r", newline="") as fcsv:
                        rdr = csv.reader(fcsv)
                        next(rdr, None)
                        for row in rdr:
                            if len(row) >= 3:
                                try:
                                    frq.append(float(row[0])); vel.append(float(row[1])); wl.append(float(row[2]))
                                except Exception:
                                    continue
                except Exception:
                    continue
                if frq:
                    combined_rows.append((offset, frq, vel, wl))
            if combined_rows and len(combined_rows) > 1:
                min_len = min(len(r[1]) for r in combined_rows)
                comb_csv = os.path.join(self.output_folder, f"combined_{key}.csv")
                with open(comb_csv, "w", newline="") as fcsv:
                    w = csv.writer(fcsv)
                    header = []
                    for off, _, _, _ in combined_rows:
                        header += [f"freq({key}_{off})", f"vel({key}_{off})", f"wav({key}_{off})"]
                    w.writerow(header)
                    for i in range(min_len):
                        row = []
                        for _, frq, vel, wl in combined_rows:
                            row += [frq[i], vel[i], wl[i] if i < len(wl) else ""]
                        w.writerow(row)
                if self.logbox:
                    self.logbox.insert(tk.END, f"Combined CSV → {comb_csv}\n"); self.logbox.see(tk.END)
        except Exception:
            pass

    def _build_combined_csv_compare(self, paths: list[str]):
        """Aggregate per-file compare CSVs into combined_compare_all.csv."""
        try:
            import csv
            rows_by_file = []
            headers = []
            for pth in paths:
                b = os.path.splitext(os.path.basename(pth))[0]
                comb_path = os.path.join(self.output_folder, f"{b}_compare.csv")
                if not os.path.isfile(comb_path):
                    continue
                with open(comb_path, "r", newline="") as fcsv:
                    rdr = list(csv.reader(fcsv))
                    if not rdr:
                        continue
                    hdr = rdr[0]; data = rdr[1:]
                    headers += [f"{h}[{b}]" for h in hdr]
                    rows_by_file.append(data)
            if rows_by_file and len(rows_by_file) > 1:
                min_len = min(len(r) for r in rows_by_file)
                out_all = os.path.join(self.output_folder, "combined_compare_all.csv")
                with open(out_all, "w", newline="") as fcsv:
                    w = csv.writer(fcsv)
                    w.writerow(headers)
                    for i in range(min_len):
                        row = []
                        for data in rows_by_file:
                            row += data[i]
                        w.writerow(row)
                if self.logbox:
                    self.logbox.insert(tk.END, f"Combined compare CSV → {out_all}\n"); self.logbox.see(tk.END)
        except Exception:
            pass


def main() -> None:
    root = tk.Tk(); SimpleMASWGUI(root); root.mainloop()


if __name__ == "__main__":
    main()


